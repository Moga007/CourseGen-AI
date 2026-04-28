"""
Module de conversion du cours Markdown en présentation PowerPoint (.pptx).
Thème académique sombre – design splitté et soigné.
"""

import re
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml import etree

# ═══════════════════════════════════════════════════════
#  PALETTE
# ═══════════════════════════════════════════════════════
C_BG           = RGBColor(0x08, 0x08, 0x18)   # Fond principal – très sombre
C_BG_CARD      = RGBColor(0x0F, 0x0F, 0x24)   # Fond slides contenu
C_PANEL        = RGBColor(0x5B, 0x5E, 0xE8)   # Panneau gauche titre (violet vif)
C_ACCENT       = RGBColor(0x63, 0x66, 0xF1)   # Accent principal
C_ACCENT2      = RGBColor(0x81, 0x8C, 0xF8)   # Accent clair
C_ACCENT_DARK  = RGBColor(0x1E, 0x20, 0x5C)   # Violet très sombre
C_ACCENT_MID   = RGBColor(0x2D, 0x30, 0x7A)   # Violet intermédiaire
C_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT_LIGHT   = RGBColor(0xE2, 0xE8, 0xF0)
C_TEXT_MUTED   = RGBColor(0x94, 0xA3, 0xB8)
C_TABLE_HDR    = RGBColor(0x3B, 0x3E, 0xB8)
C_TABLE_ROW1   = RGBColor(0x12, 0x12, 0x28)
C_TABLE_ROW2   = RGBColor(0x18, 0x18, 0x34)
C_SUCCESS      = RGBColor(0x34, 0xD3, 0x99)   # vert (pour / avantages / après)
C_WARN         = RGBColor(0xF5, 0x9E, 0x42)   # orange (contre / inconvénients / avant)
C_BULLET       = RGBColor(0x6E, 0x75, 0xF9)   # Couleur des marqueurs bullets

# ═══════════════════════════════════════════════════════
#  TYPOGRAPHIE
# ═══════════════════════════════════════════════════════
# Aptos est la police moderne par défaut de Microsoft 365 (2023+).
# Aptos Display est optimisée pour les grands titres (contraste + lisibilité).
# Fallback automatique vers Calibri si Aptos indisponible.
FONT_DISPLAY = 'Aptos Display'   # Titres, grands chiffres, labels de section
FONT_BODY    = 'Aptos'           # Corps de texte, puces, descriptions
FONT_MONO    = 'Consolas'        # Code source (Consolas est sur tout Windows)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

TITLE_PANEL_W  = Inches(4.4)    # Largeur du panneau gauche (slide titre)
LEFT_BAR_W     = Inches(0.09)   # Barre gauche accent sur slides contenu

MAX_BULLETS_SINGLE = 5
MAX_BULLETS_COL    = 7
MAX_CHARS          = 400

CHIFFRES_ROMAINS = ['I', 'II', 'III', 'IV', 'V', 'VI', 'VII', 'VIII', 'IX', 'X']
CIRCLED_NUMS     = ['①', '②', '③', '④', '⑤', '⑥', '⑦', '⑧', '⑨', '⑩',
                    '⑪', '⑫']

# Chips pédagogiques par type de slide (médaillons circulaires à côté du titre)
CHIP_DEF     = '◈'   # définitions de concepts
CHIP_PTS     = '★'   # points importants à retenir
CHIP_STATS   = '%'   # statistiques / chiffres clés
CHIP_SCHEMA  = '◉'   # schéma / relations entre éléments
CHIP_COMPARE = '⇌'   # comparaison (deux colonnes)
CHIP_TABLE   = '▦'   # tableau de données
CHIP_STEPS   = '⇣'   # processus séquentiel / stepper vertical
CHIP_GOAL    = '◎'   # objectifs pédagogiques (cible / visée)
CHIP_RECAP   = '✦'   # synthèse / à retenir (distinct de CHIP_PTS='★')
CHIP_CHECK   = '✓'   # marqueur de validation devant chaque objectif
CHIP_TIMELINE = '→'  # frise chronologique / timeline (évoque la progression)
CHIP_CASE    = '▣'   # cas pratique / étude de cas (frame le sujet étudié)
CHIP_CODE    = '⌨'   # bloc de code (clavier — universellement reconnu)

# Marque affichée dans le pied de page (badge gauche). À adapter par école.
BRAND_LABEL = 'IESIG'

# Motif de fond subtil (appliqué à toutes les slides de contenu via _content_base).
# Presets DrawingML les plus adaptés à un fond sombre :
#   'smGrid'    : grille fine (aspect technique/académique)
#   'dotGrid'   : grille de points (le plus discret)
#   'ltDnDiag'  : diagonales fines descendantes (aspect « papier »)
#   'trellis'   : motif tressé léger
# Alpha en % : 3-6 pour rester quasi-invisible sans écraser le texte.
BG_PATTERN_ENABLED = True
BG_PATTERN_PRESET  = 'dotGrid'
BG_PATTERN_COLOR   = RGBColor(0x81, 0x8C, 0xF8)  # C_ACCENT2 (clair → plus visible sur fond sombre)
BG_PATTERN_ALPHA   = 5    # %


# ═══════════════════════════════════════════════════════
#  UTILITAIRES MARKDOWN
# ═══════════════════════════════════════════════════════

def _clean(text: str) -> str:
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*',     r'\1', text)
    text = re.sub(r'`(.+?)`',       r'\1', text)
    text = re.sub(r'__(.+?)__',     r'\1', text)
    return text.strip()


def _truncate(text: str, max_chars: int = MAX_CHARS) -> str:
    text = _clean(text)
    if len(text) <= max_chars:
        return text
    return text[:max_chars].rsplit(' ', 1)[0] + '…'


def _parse_inline(text: str) -> list[tuple[str, bool]]:
    """Retourne (texte, est_gras) pour préserver le gras Markdown."""
    segments, last = [], 0
    for m in re.finditer(r'\*\*(.+?)\*\*|__(.+?)__', text):
        if m.start() > last:
            plain = re.sub(r'\*(.+?)\*', r'\1', text[last:m.start()])
            plain = re.sub(r'`(.+?)`',   r'\1', plain)
            if plain:
                segments.append((plain, False))
        segments.append((_clean(m.group(1) or m.group(2)), True))
        last = m.end()
    if last < len(text):
        tail = re.sub(r'\*(.+?)\*', r'\1', text[last:])
        tail = re.sub(r'`(.+?)`',   r'\1', tail)
        if tail.strip():
            segments.append((tail, False))
    return segments or [(_clean(text), False)]


def _extract_bullets(text: str, max_b: int = MAX_BULLETS_COL * 2) -> list[str]:
    bullets = []
    for line in text.split('\n'):
        line = line.strip()
        if re.match(r'^[-*•]\s', line):
            b = line[2:].strip()
            if b:
                bullets.append(b)
        elif re.match(r'^\d+\.\s', line):
            b = re.sub(r'^\d+\.\s', '', line).strip()
            if b:
                bullets.append(b)
        if len(bullets) >= max_b:
            break
    return bullets


def _first_paragraph(text: str) -> str:
    for block in text.split('\n\n'):
        block = block.strip()
        if block and not block.startswith('|') and not block.startswith('#'):
            return _truncate(block)
    return ''


def _wrap_text(text: str, width: int = 90) -> list[str]:
    words = text.split()
    lines, current = [], ''
    for word in words:
        if len(current) + len(word) + 1 <= width:
            current = (current + ' ' + word).strip()
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines or [text]


def _parse_md_table(text: str) -> tuple[list[str], list[list[str]]]:
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    tbl = [l for l in lines if l.startswith('|')]
    if len(tbl) < 2:
        return [], []
    def parse_row(line):
        return [_clean(c.strip()) for c in line.strip('|').split('|')]
    headers = parse_row(tbl[0])
    rows = [parse_row(l) for l in tbl[2:] if not re.match(r'^\|[-:| ]+\|$', l)]
    return headers, rows


def _parse_definitions(text: str) -> list[tuple[str, str]]:
    items = []
    for line in text.split('\n'):
        line = line.strip().lstrip('-').lstrip('*').strip()
        m = re.match(r'\*\*(.+?)\*\*\s*[:\-–]\s*(.+)', line)
        if m:
            items.append((_clean(m.group(1)), m.group(2).strip()))
            continue
        m2 = re.match(r'^([^:\-]{3,40})\s*[:\-–]\s*(.{10,})', line)
        if m2:
            items.append((_clean(m2.group(1)), m2.group(2).strip()))
    return items[:12]


def _parse_blockquote(text: str) -> tuple[str, str]:
    """
    Extrait le corps et l'attribution d'un blockquote markdown.
    Accepte des formes comme :

        > La pensée devient réalité par l'action.
        > — Albert Einstein

    La dernière ligne commençant par « — » ou « - » (tiret) devient
    l'attribution, le reste est le corps de la citation.
    Retourne ('', '') si aucun blockquote n'est trouvé.
    """
    quote_lines = []
    for raw in text.split('\n'):
        line = raw.strip()
        if line.startswith('>'):
            # Enlève le marqueur '>' et les espaces qui suivent
            quote_lines.append(line.lstrip('>').strip())
        elif quote_lines and not line:
            # Ligne vide après des quotes : on continue, elles peuvent être
            # entrecoupées de lignes vides (markdown standard)
            continue
        elif quote_lines:
            # Ligne non-quote rencontrée après des quotes : on s'arrête
            break

    # Enlève les lignes vides éventuelles conservées
    quote_lines = [l for l in quote_lines if l]
    if not quote_lines:
        return '', ''

    attribution = ''
    # Détecte une attribution en dernière ligne (« — Auteur », « - Source »)
    last = quote_lines[-1]
    m = re.match(r'^[—–\-]\s*(.+)', last)
    if m and len(quote_lines) > 1:
        attribution = m.group(1).strip()
        quote_lines = quote_lines[:-1]

    body = ' '.join(quote_lines).strip()
    return _clean(body), _clean(attribution)


def _is_pure_blockquote(text: str) -> bool:
    """True si toutes les lignes non vides du texte sont des quotes markdown."""
    has_quote = False
    for line in text.split('\n'):
        line = line.strip()
        if not line:
            continue
        if line.startswith('>'):
            has_quote = True
        else:
            return False
    return has_quote


def _parse_steps(text: str) -> list[dict]:
    """
    Parse une liste numérotée markdown en étapes structurées (stepper).

    Formats supportés :
        1. **Titre de l'étape** : description...
        2. **Titre** : description multi-lignes
           qui continue ici.
        3. Titre sans gras — description éventuelle
        4. Ligne unique

    Retourne [{'titre': str, 'description': str}, ...] (dans l'ordre du doc).
    """
    steps = []
    current = None
    for raw in text.split('\n'):
        line = raw.strip()
        m = re.match(r'^\d+\.\s+(.*)', line)
        if m:
            if current:
                steps.append(current)
            body = m.group(1).strip()
            # Pattern principal : **Titre** : description
            mb = re.match(r'^\*\*(.+?)\*\*\s*[:：\-–]?\s*(.*)', body)
            if mb:
                titre = _clean(mb.group(1))
                desc  = _clean(mb.group(2)) if mb.group(2) else ''
            elif ':' in body:
                tit, desc = body.split(':', 1)
                titre, desc = _clean(tit), _clean(desc)
            elif ' — ' in body or ' – ' in body:
                sep = ' — ' if ' — ' in body else ' – '
                tit, desc = body.split(sep, 1)
                titre, desc = _clean(tit), _clean(desc)
            else:
                titre, desc = _clean(body), ''
            current = {'titre': titre, 'description': desc}
        elif current and line and not line.startswith('#') and not line.startswith('-'):
            # Continuation de la description précédente
            current['description'] = (current['description'] + ' ' + _clean(line)).strip()
    if current:
        steps.append(current)
    return steps


_KPI_VALUE_RE = re.compile(
    r'^\s*[-*•]\s*\*\*([^*]+?)\*\*\s*(?:[:：\-–—]|\s)\s*(.+)$'
    r'|^\s*[-*•]\s*(\d[\d\s.,]*\s*(?:%|€|\$|£|¥|M|Md|Mds|k|×|x|pts?|j|h|mois|ans?|fois)?)\s*'
    r'(?:[:：\-–—])\s*(.+)$',
    re.IGNORECASE,
)

_KPI_LOOKS_NUMERIC_RE = re.compile(
    r'^\s*\d[\d\s.,]*\s*(%|€|\$|£|¥|M|Md|Mds|k|×|x|pts?|j|h|mois|ans?|fois)?\s*$',
    re.IGNORECASE,
)


def _parse_kpi_bullets(text: str, max_n: int = 4) -> list[dict]:
    """
    Extrait des KPI chiffrés depuis des bullets markdown.

    Patterns supportés (majoritairement rencontrés dans les cours IESIG) :
        - **85%** des entreprises B2B utilisent le CRM
        - **1 500 €** par client et par an
        - **3×** plus efficace que la méthode traditionnelle
        - 73% — des français consultent les avis en ligne

    Retourne [{'valeur': str, 'label': str}, ...] (max_n items).
    La valeur doit être purement numérique (avec suffixe éventuel) pour être
    reconnue comme KPI. Le label est la suite du bullet, tronqué à 60 chars.
    """
    out = []
    for raw in text.split('\n'):
        line = raw.strip()
        if not line or not re.match(r'^[-*•]', line):
            continue
        m = _KPI_VALUE_RE.match(line)
        if not m:
            continue
        # Deux captures possibles selon l'alternative du regex
        if m.group(1) is not None:
            valeur_candidate = m.group(1).strip()
            label_candidate  = m.group(2).strip()
        else:
            valeur_candidate = m.group(3).strip()
            label_candidate  = m.group(4).strip()

        # On ne garde que si la valeur est effectivement numérique
        if not _KPI_LOOKS_NUMERIC_RE.match(valeur_candidate):
            continue
        valeur = _clean(valeur_candidate)
        label  = _truncate(_clean(label_candidate), 60)
        if valeur and label:
            out.append({'valeur': valeur, 'label': label})
        if len(out) >= max_n:
            break
    return out


# Token date pour timeline : 4 chiffres (1995), 1995s, 1990-2000, "années 90",
# mois français + année, Q1 2020.
_TIMELINE_DATE_TOKEN = (
    r'(?:\d{4}(?:s|\s*[-–]\s*\d{2,4})?'
    r'|années?\s+\d{2,4}'
    r'|(?:janvier|février|fevrier|mars|avril|mai|juin|juillet|août|aout|'
    r'septembre|octobre|novembre|décembre|decembre)\s+\d{4}'
    r'|q[1-4]\s+\d{4})'
)

# Cas A : - **2020 : Titre** : description (date + titre dans le gras)
_TIMELINE_RE_COMBO = re.compile(
    r'^\s*[-*•]\s*\*\*\s*(' + _TIMELINE_DATE_TOKEN + r'[^*]*?)\s*\*\*'
    r'\s*[:：\-–—]?\s*(.*)$',
    re.IGNORECASE,
)
# Cas B : - **2020** : titre — description  /  - 2020 — titre : description
_TIMELINE_RE_SPLIT = re.compile(
    r'^\s*[-*•]\s*(?:\*\*)?(' + _TIMELINE_DATE_TOKEN + r')(?:\*\*)?'
    r'\s*[:：\-–—]\s*(.+)$',
    re.IGNORECASE,
)


def _parse_timeline(text: str, max_n: int = 7) -> list[dict]:
    """
    Parse une liste d'événements datés (frise chronologique).

    Patterns supportés :
        - **1995** : Création de Yahoo!
        - 1995 — Création de Yahoo!
        - **2010 : Avènement du mobile** : 50% du trafic web
        - 1990s — Décennie de l'expansion
        - **Années 90** : massification d'Internet
        - **Janvier 2020** : COVID et bascule digitale

    Retourne [{'date': str, 'titre': str, 'description': str}, ...].
    """
    out = []
    for raw in text.split('\n'):
        line = raw.strip()
        if not line:
            continue

        # Cas A : ** englobe date + titre
        mc = _TIMELINE_RE_COMBO.match(line)
        if mc:
            combo = _clean(mc.group(1))
            # Sépare la date du reste : "2020 : Titre" ou "2020 — Titre"
            ms = re.match(
                r'^(' + _TIMELINE_DATE_TOKEN + r')\s*[:：\-–—]?\s*(.*)$',
                combo, re.IGNORECASE,
            )
            if ms:
                date  = ms.group(1).strip()
                titre = ms.group(2).strip() or ''
            else:
                date  = combo
                titre = ''
            desc = _clean(mc.group(2))
            if not titre and desc:
                # Si pas de titre interne, le reste après ** devient le titre
                titre, desc = desc, ''
            out.append({'date': date, 'titre': titre, 'description': desc})
            if len(out) >= max_n:
                break
            continue

        # Cas B : date au début, séparateur, suite de la ligne
        ms = _TIMELINE_RE_SPLIT.match(line)
        if ms:
            date = _clean(ms.group(1))
            rest = _clean(ms.group(2))
            # rest peut contenir un sous-titre gras
            mt = re.match(r'^\*\*(.+?)\*\*\s*[:：\-–—]?\s*(.*)', rest)
            if mt:
                titre = _clean(mt.group(1))
                desc  = _clean(mt.group(2))
            else:
                # Sépare titre/description sur " : " ou " — "
                titre, desc = rest, ''
                for sep in (' : ', ' — ', ' – '):
                    if sep in rest:
                        titre, desc = rest.split(sep, 1)
                        titre, desc = _clean(titre), _clean(desc)
                        break
            out.append({'date': date, 'titre': titre, 'description': desc})
            if len(out) >= max_n:
                break
    return out


def _looks_like_timeline(text: str) -> bool:
    """
    True si >=3 bullets et au moins 60% sont des événements datés
    (heuristique pour auto-détecter une frise sans mot-clé dans le titre).
    """
    bullets = [l for l in text.split('\n') if re.match(r'^\s*[-*•]', l.strip())]
    if len(bullets) < 3:
        return False
    events = _parse_timeline(text, max_n=20)
    return len(events) >= 3 and len(events) >= 0.6 * len(bullets)


# Synonymes par rôle dans un cas pratique (utilisés pour classifier les
# H3/labels en l'une des 4 cases : contexte / problème / solution / résultat).
_CASE_LABELS = {
    'contexte': ('contexte', 'context', 'situation', 'cadre', 'arrière-plan',
                 'arriere-plan', 'mise en contexte', 'cas étudié', 'cas etudie',
                 'environnement'),
    'probleme': ('problème', 'probleme', 'problem', 'enjeu', 'défi', 'defi',
                 'difficulté', 'difficulte', 'challenge', 'question',
                 'questionnement', 'point bloquant'),
    'solution': ('solution', 'approche', 'réponse', 'reponse', 'démarche',
                 'demarche', 'action', 'actions', 'stratégie', 'strategie',
                 'mise en œuvre', 'mise en oeuvre', 'plan d\'action'),
    'resultat': ('résultat', 'resultat', 'résultats', 'resultats', 'outcome',
                 'impact', 'bilan', 'effets', 'effet', 'résolution',
                 'resolution', 'leçons', 'lecons', 'enseignements',
                 'retour d\'expérience', 'retour dexperience',
                 'conclusion du cas'),
}


def _classify_case_label(label: str) -> str | None:
    """Classifie un label en l'une des 4 cases du cas pratique."""
    lbl = label.lower().strip().rstrip(':：').strip()
    for key, kws in _CASE_LABELS.items():
        for kw in kws:
            if kw in lbl:
                return key
    return None


def _parse_case_study(text: str) -> dict | None:
    """
    Parse un cas pratique en 4 cases : contexte, problème, solution, résultat.

    Patterns supportés (testés dans cet ordre) :
      1. Sous-titres H3 / H4 :
            ### Contexte
            <body>
            ### Problème
            <body>
            ### Solution
            <body>
            ### Résultat
            <body>

      2. Inline gras :
            **Contexte** : <body>
            **Problème** : <body>
            ...

      3. Bullets gras :
            - **Contexte** : <body>
            - **Problème** : <body>
            ...

    Retourne {'contexte', 'probleme', 'solution', 'resultat'} dont chaque
    valeur peut être '' si non détectée. Renvoie None si moins de 3 cases
    sur 4 sont remplies (un cas pratique sans 3/4 c'est un faux positif).
    """
    out = {'contexte': '', 'probleme': '', 'solution': '', 'resultat': ''}

    # Cas 1 : H3/H4
    h3_re = re.compile(
        r'^#{3,4}\s+(.+?)\s*$\n?(.*?)(?=^#{3,4}\s|\Z)',
        re.MULTILINE | re.DOTALL,
    )
    for m in h3_re.finditer(text):
        slot = _classify_case_label(m.group(1))
        if slot and not out[slot]:
            out[slot] = _clean(m.group(2)).strip()

    # Cas 2 : inline **Label** : ... (sépare sur \n\n ou prochain **)
    if not all(out.values()):
        inline_re = re.compile(
            r'\*\*\s*([^*\n:]+?)\s*\*\*\s*[:：]\s*(.+?)(?=\n\s*\*\*|\n\n|\Z)',
            re.DOTALL,
        )
        for m in inline_re.finditer(text):
            slot = _classify_case_label(m.group(1))
            if slot and not out[slot]:
                out[slot] = _clean(m.group(2)).strip()

    # Cas 3 : bullets - **Label** : ...
    if not all(out.values()):
        for raw in text.split('\n'):
            line = raw.strip()
            mb = re.match(r'^[-*•]\s*\*\*([^*]+?)\*\*\s*[:：]\s*(.+)$', line)
            if mb:
                slot = _classify_case_label(mb.group(1))
                if slot and not out[slot]:
                    out[slot] = _clean(mb.group(2)).strip()

    if sum(1 for v in out.values() if v) < 3:
        return None
    return out


# Préfixes de commentaire par langage (pour mise en couleur muted automatique)
_CODE_COMMENT_PREFIXES = {
    'python':   ('#',),       'py':   ('#',),
    'bash':     ('#',),       'sh':   ('#',),
    'ruby':     ('#',),       'rb':   ('#',),
    'r':        ('#',),
    'yaml':     ('#',),       'yml':  ('#',),
    'toml':     ('#',),
    'javascript': ('//',),    'js':   ('//',),
    'typescript': ('//',),    'ts':   ('//',),
    'jsx':      ('//',),      'tsx':  ('//',),
    'java':     ('//',),
    'c':        ('//',),      'cpp':  ('//',),
    'csharp':   ('//',),      'cs':   ('//',),
    'go':       ('//',),
    'rust':     ('//',),      'rs':   ('//',),
    'php':      ('//',),
    'swift':    ('//',),
    'kotlin':   ('//',),      'kt':   ('//',),
    'scala':    ('//',),
    'sql':      ('--',),
    'html':     ('<!--',),    'xml':  ('<!--',),
    'css':      ('/*',),
}


def _is_comment_line(line: str, language: str) -> bool:
    """True si la ligne de code est un commentaire selon le langage."""
    prefixes = _CODE_COMMENT_PREFIXES.get(language.lower(), ())
    if not prefixes:
        return False
    stripped = line.lstrip()
    return any(stripped.startswith(p) for p in prefixes)


def _parse_code_blocks(text: str):
    """
    Extrait le PREMIER bloc de code markdown ```lang ... ``` du texte.
    Retourne {'language', 'code', 'before', 'after'} ou None si pas de bloc.

    Le langage est optionnel (forme ``` ... ``` sans étiquette acceptée).
    'before' = texte avant le bloc (utilisé comme description),
    'after'  = texte après le bloc (rarement utilisé pour l'instant).
    """
    pattern = re.compile(
        r'^```([a-zA-Z0-9_+\-]*)\s*\n(.*?)\n```\s*$',
        re.MULTILINE | re.DOTALL,
    )
    m = pattern.search(text)
    if not m:
        return None
    language = m.group(1).strip()
    code = m.group(2)
    before = text[:m.start()].strip()
    after = text[m.end():].strip()
    if not code.strip():
        return None
    return {
        'language': language,
        'code': code,
        'before': before,
        'after': after,
    }


def _looks_like_kpi(text: str) -> bool:
    """
    True si le texte ressemble à une slide KPI : au moins 2 bullets et
    >=50% des bullets reconnus comme KPI par `_parse_kpi_bullets`.
    """
    bullets = [l for l in text.split('\n') if re.match(r'^\s*[-*•]', l.strip())]
    if len(bullets) < 2:
        return False
    kpis = _parse_kpi_bullets(text, max_n=12)
    return len(kpis) >= 2 and len(kpis) >= 0.5 * len(bullets)


def _looks_like_steps(text: str) -> bool:
    """
    Heuristique : True si le contenu est une liste numérotée de 3-8 étapes
    qui domine le texte (pour auto-détection stepper).
    """
    lines = [l.strip() for l in text.split('\n') if l.strip()]
    if not lines:
        return False
    numbered = sum(1 for l in lines if re.match(r'^\d+\.\s', l))
    if not (3 <= numbered <= 8):
        return False
    # Première ligne non vide doit être une étape numérotée (pas un paragraphe d'intro)
    if not re.match(r'^\d+\.\s', lines[0]):
        return False
    return True


# ═══════════════════════════════════════════════════════
#  HELPERS PPTX
# ═══════════════════════════════════════════════════════

def _fill_solid(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def _rect(slide, left, top, width, height, color: RGBColor):
    """Ajoute un rectangle plein sans bord."""
    s = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    _fill_solid(s, color)
    _no_line(s)
    return s


def _rounded_rect(slide, left, top, width, height, color: RGBColor):
    """Rectangle arrondi plein sans bord (type 5)."""
    s = slide.shapes.add_shape(5, Inches(left), Inches(top), Inches(width), Inches(height))
    _fill_solid(s, color)
    _no_line(s)
    return s


def _tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )


def _enable_shrink_to_fit(tf):
    """
    Active la réduction automatique de la police et de l'interligne quand
    le texte déborde du cadre (équivalent au « Réduire le texte en cas de
    dépassement » de PowerPoint, cf. <a:normAutofit/>).

    PowerPoint calcule lui-même fontScale et lnSpcReduction au rendu :
    si tout tient → aucun changement ; sinon les polices sont réduites
    proportionnellement (la hiérarchie typographique est préservée).

    Précondition : le text frame doit avoir une hauteur fixe (ce que fait
    toujours `_tb`). Sans hauteur bornée, la réduction n'a pas de cible.

    Idempotent : retire d'abord les éventuels éléments noAutofit /
    spAutoFit / normAutofit existants pour éviter les doublons.
    """
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is None:
        return
    for tag in ('a:noAutofit', 'a:spAutoFit', 'a:normAutofit'):
        for el in bodyPr.findall(qn(tag)):
            bodyPr.remove(el)
    etree.SubElement(bodyPr, qn('a:normAutofit'))


def _set_bg(slide, color: RGBColor):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _run_fmt(run, size_pt: int, color: RGBColor, bold=False, italic=False,
             font: str = FONT_BODY):
    run.font.size    = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.name    = font


def _add_runs(para, raw_text: str, size_pt: int, color: RGBColor, base_bold=False,
              font: str = FONT_BODY):
    for seg_text, is_bold in _parse_inline(raw_text):
        run = para.add_run()
        run.text = seg_text
        _run_fmt(run, size_pt, color, bold=(base_bold or is_bold), font=font)


def _add_transparent_rect(slide, left, top, width, height, color: RGBColor, opacity: int):
    """Rectangle avec transparence. opacity: 0=invisible, 100=opaque (%).

    On passe par l'API haut niveau (`_fill_solid`) pour que le solidFill créé
    surcharge correctement le `p:style` par défaut du shape (qui pointe vers
    le thème Office — fill bleu clair d'accent 2). Sans ça, le style thème
    prend le dessus et la couleur/alpha ne sont pas appliqués.
    Puis on injecte l'élément <a:alpha> dans le srgbClr déjà créé.
    """
    s = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    _no_line(s)
    _fill_solid(s, color)

    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    spPr = s._element.find(qn('p:spPr'))
    if spPr is None:
        return s
    solid = spPr.find(qn('a:solidFill'))
    if solid is None:
        return s
    srgb = solid.find(qn('a:srgbClr'))
    if srgb is None:
        return s
    # Nettoie un éventuel alpha existant pour éviter les doublons
    for ex in srgb.findall(qn('a:alpha')):
        srgb.remove(ex)
    alpha = etree.SubElement(srgb, f'{{{ns}}}alpha')
    # DrawingML ST_PositivePercentage : 100000 = 100 % opaque
    alpha.set('val', str(int(opacity * 1000)))
    return s


_DML_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def _add_pattern_overlay(slide, preset: str = None, fg_color: RGBColor = None,
                         fg_alpha_pct: int = None):
    """
    Motif de fond subtil pleine-slide (pattFill DrawingML natif).

    Crée un rectangle pleine-slide rempli d'un motif preset (dotGrid, smGrid,
    ltDnDiag, etc.) avec fond transparent et avant-plan très peu opaque pour
    donner une texture discrète sans écraser le contenu. La forme est posée
    en premier donc tous les appels suivants de _content_base passent au-dessus.
    """
    preset = preset or BG_PATTERN_PRESET
    fg_color = fg_color or BG_PATTERN_COLOR
    alpha = BG_PATTERN_ALPHA if fg_alpha_pct is None else fg_alpha_pct

    s = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    _no_line(s)

    sp = s._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return s
    # Supprimer tout fill existant pour y poser notre pattFill
    for child in list(spPr):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('solidFill', 'gradFill', 'pattFill', 'noFill', 'blipFill'):
            spPr.remove(child)
    # Neutraliser le p:style du shape (sinon fillRef du thème surcharge notre pattFill)
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)

    patt = etree.SubElement(spPr, f'{{{_DML_NS}}}pattFill')
    patt.set('prst', preset)
    fg = etree.SubElement(patt, f'{{{_DML_NS}}}fgClr')
    fg_srgb = etree.SubElement(fg, f'{{{_DML_NS}}}srgbClr')
    fg_srgb.set('val', f'{fg_color[0]:02X}{fg_color[1]:02X}{fg_color[2]:02X}')
    fg_alpha = etree.SubElement(fg_srgb, f'{{{_DML_NS}}}alpha')
    fg_alpha.set('val', str(int(alpha * 1000)))
    bg = etree.SubElement(patt, f'{{{_DML_NS}}}bgClr')
    bg_srgb = etree.SubElement(bg, f'{{{_DML_NS}}}srgbClr')
    bg_srgb.set('val', '000000')
    bg_alpha = etree.SubElement(bg_srgb, f'{{{_DML_NS}}}alpha')
    bg_alpha.set('val', '0')  # fond du pattern totalement transparent
    return s


def _add_slide_hyperlink(run, source_slide, target_slide):
    """
    Ajoute un hyperlien interne (saut vers une autre slide) sur un run de texte.

    Crée (ou réutilise) une relation de type 'slide' entre la slide source et
    la slide cible, puis injecte un <a:hlinkClick> dans les propriétés du run
    avec l'action ppaction://hlinksldjump qui indique à PowerPoint d'exécuter
    un saut de slide au clic.
    """
    if target_slide is None or source_slide is None:
        return
    rId = source_slide.part.relate_to(target_slide.part, RT.SLIDE)

    rPr = run._r.get_or_add_rPr()
    # Supprime un hlinkClick antérieur éventuel (évite les doublons)
    for existing in rPr.findall(qn('a:hlinkClick')):
        rPr.remove(existing)
    hlink = etree.SubElement(rPr, qn('a:hlinkClick'))
    hlink.set(qn('r:id'), rId)
    hlink.set('action', 'ppaction://hlinksldjump')


def _fill_gradient(shape, stops: list, angle_deg: float = 90.0):
    """
    Applique un gradient linéaire. stops = [(pos_0_100, RGBColor), ...].
    angle_deg : 0 = gauche→droite, 90 = haut→bas, 45 = diagonal descendant.
    """
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    for child in list(spPr):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('solidFill', 'gradFill', 'pattFill', 'noFill', 'blipFill'):
            spPr.remove(child)
    grad = etree.SubElement(spPr, f'{{{_DML_NS}}}gradFill')
    grad.set('flip', 'none')
    grad.set('rotWithShape', '1')
    gsLst = etree.SubElement(grad, f'{{{_DML_NS}}}gsLst')
    for pos_pct, col in stops:
        gs = etree.SubElement(gsLst, f'{{{_DML_NS}}}gs')
        gs.set('pos', str(int(pos_pct * 1000)))
        srgb = etree.SubElement(gs, f'{{{_DML_NS}}}srgbClr')
        srgb.set('val', f'{col[0]:02X}{col[1]:02X}{col[2]:02X}')
    lin = etree.SubElement(grad, f'{{{_DML_NS}}}lin')
    lin.set('ang', str(int(angle_deg * 60000)))
    lin.set('scaled', '0')
    shape.line.fill.background()
    return shape


def _add_shadow(shape, blur_pt: int = 8, dist_pt: int = 4,
                alpha_pct: int = 40, dir_deg: int = 90):
    """
    Ombre portée subtile. blur_pt/dist_pt en points, alpha_pct 0-100,
    dir_deg : 90 = vers le bas, 135 = bas-droite.
    """
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    for child in list(spPr):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag == 'effectLst':
            spPr.remove(child)
    effectLst = etree.Element(f'{{{_DML_NS}}}effectLst')
    outerShdw = etree.SubElement(effectLst, f'{{{_DML_NS}}}outerShdw')
    outerShdw.set('blurRad', str(blur_pt * 12700))
    outerShdw.set('dist', str(dist_pt * 12700))
    outerShdw.set('dir', str(int(dir_deg * 60000)))
    outerShdw.set('algn', 'ctr')
    outerShdw.set('rotWithShape', '0')
    srgb = etree.SubElement(outerShdw, f'{{{_DML_NS}}}srgbClr')
    srgb.set('val', '000000')
    alpha = etree.SubElement(srgb, f'{{{_DML_NS}}}alpha')
    alpha.set('val', str(alpha_pct * 1000))
    # Position correcte dans l'ordre des enfants de spPr : après <a:ln> si présent.
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        ln.addnext(effectLst)
    else:
        spPr.append(effectLst)
    return shape


def _add_footer(slide, brand: str, breadcrumb: str, page_num: int, total: int):
    """
    Pied de page institutionnel à 3 zones :

      [BADGE]   spécialité › module › chapitre                    07 / 28
      ─────────────────────────────────────────────────────────────────────

    - Badge gauche : marque (IESIG par défaut), pill arrondi accent
    - Centre : breadcrumb hiérarchique avec séparateurs U+203A
    - Droite : numéro de page courant en gras accent + total muted
    - Trait fin C_ACCENT_MID au-dessus, sur toute la largeur utile
    """
    footer_y = 7.28
    footer_h = 0.20

    # Trait séparateur très fin au-dessus du footer
    _rect(slide, 0.35, footer_y - 0.10, 12.63, 0.012, C_ACCENT_MID)

    # ── Badge marque (gauche) ────────────────────────────────────────
    if brand:
        bw = max(0.55, len(brand) * 0.085 + 0.20)
        pill = _rounded_rect(slide, 0.35, footer_y, bw, footer_h, C_ACCENT)
        _add_shadow(pill, blur_pt=4, dist_pt=1, alpha_pct=30, dir_deg=90)
        tb_b = _tb(slide, 0.35, footer_y, bw, footer_h)
        tf_b = tb_b.text_frame
        tf_b.margin_left   = Inches(0.04)
        tf_b.margin_right  = Inches(0.04)
        tf_b.margin_top    = Inches(0)
        tf_b.margin_bottom = Inches(0)
        tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_b.word_wrap = False
        p_b = tf_b.paragraphs[0]
        p_b.alignment = PP_ALIGN.CENTER
        run_b = p_b.add_run()
        run_b.text = brand
        _run_fmt(run_b, 8, C_WHITE, bold=True, font=FONT_DISPLAY)
        crumb_left = 0.35 + bw + 0.20
    else:
        crumb_left = 0.35

    # ── Breadcrumb (centre, aligné gauche après le badge) ────────────
    crumb_w = 9.3 - crumb_left
    if crumb_w > 0.5 and breadcrumb:
        tb_c = _tb(slide, crumb_left, footer_y, crumb_w, footer_h)
        tf_c = tb_c.text_frame
        tf_c.margin_left   = Inches(0.02)
        tf_c.margin_right  = Inches(0.02)
        tf_c.margin_top    = Inches(0)
        tf_c.margin_bottom = Inches(0)
        tf_c.word_wrap = False
        tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_c = tf_c.paragraphs[0]
        p_c.alignment = PP_ALIGN.LEFT
        run_c = p_c.add_run()
        run_c.text = breadcrumb
        _run_fmt(run_c, 9, C_TEXT_MUTED, font=FONT_BODY)

    # ── Pagination (droite) ──────────────────────────────────────────
    # Numéro courant en gras accent, séparateur "/" et total en muted.
    # Cette hiérarchie typographique met l'accent sur "où on en est".
    tb_r = _tb(slide, 9.7, footer_y, 3.28, footer_h)
    tf_r = tb_r.text_frame
    tf_r.margin_left   = Inches(0)
    tf_r.margin_right  = Inches(0.04)
    tf_r.margin_top    = Inches(0)
    tf_r.margin_bottom = Inches(0)
    tf_r.word_wrap = False
    tf_r.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_r = tf_r.paragraphs[0]
    p_r.alignment = PP_ALIGN.RIGHT
    run_pn = p_r.add_run()
    run_pn.text = str(page_num)
    _run_fmt(run_pn, 11, C_ACCENT2, bold=True, font=FONT_DISPLAY)
    run_sep = p_r.add_run()
    run_sep.text = ' / '
    _run_fmt(run_sep, 9, C_TEXT_MUTED, font=FONT_BODY)
    run_tot = p_r.add_run()
    run_tot.text = str(total)
    _run_fmt(run_tot, 9, C_TEXT_MUTED, font=FONT_BODY)


def _apply_footers(prs, specialite: str, module: str, chapitre: str,
                   skip_indices: set | None = None):
    """
    Ajoute le pied de page institutionnel à toutes les slides de contenu.

    Skip systématique :
      - slide 0 (titre, full-bleed avec composition dédiée)
      - indices passés dans `skip_indices` (slides de section, qui ont
        leur propre bande basse C_ACCENT_DARK qui jurerait avec le footer)
    """
    if skip_indices is None:
        skip_indices = set()

    slides_list = list(prs.slides)
    total = len(slides_list)

    # Breadcrumb : spécialité › module › chapitre (› = U+203A)
    parts = [p.strip() for p in (specialite, module, chapitre) if p and str(p).strip()]
    breadcrumb = '  ›  '.join(parts)
    if len(breadcrumb) > 105:
        breadcrumb = breadcrumb[:102] + '…'

    for i, slide in enumerate(slides_list):
        if i == 0:
            continue
        if i in skip_indices:
            continue
        _add_footer(slide, BRAND_LABEL, breadcrumb, i + 1, total)


def _icon_chip(slide, left: float, top: float, size: float, glyph: str,
               bg_color: RGBColor = None, fg_color: RGBColor = None,
               glyph_size_pt: int = None, with_shadow: bool = True):
    """
    Médaillon circulaire avec un pictogramme centré.
    Utilisé pour marquer visuellement le type d'une slide (définition, point clé, etc.).
    """
    if bg_color is None:
        bg_color = C_ACCENT
    if fg_color is None:
        fg_color = C_WHITE
    if glyph_size_pt is None:
        glyph_size_pt = max(12, int(size * 30))

    chip = slide.shapes.add_shape(9, Inches(left), Inches(top),
                                  Inches(size), Inches(size))
    _fill_solid(chip, bg_color)
    _no_line(chip)
    if with_shadow:
        _add_shadow(chip, blur_pt=6, dist_pt=2, alpha_pct=40, dir_deg=90)

    tb = _tb(slide, left, top, size, size)
    tf = tb.text_frame
    tf.margin_left   = Inches(0.02)
    tf.margin_right  = Inches(0.02)
    tf.margin_top    = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = glyph
    _run_fmt(run, glyph_size_pt, fg_color, bold=True, font=FONT_DISPLAY)
    return chip


def _section_badge(slide, label: str, left: float, top: float):
    """Badge pill arrondi pour l'étiquette de section."""
    badge_w = min(len(label) * 0.085 + 0.4, 5.0)
    pill = _rounded_rect(slide, left, top, badge_w, 0.28, C_ACCENT_MID)
    tb = _tb(slide, left + 0.1, top + 0.02, badge_w - 0.1, 0.25)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label.upper()
    _run_fmt(run, 8, C_ACCENT2, bold=True, font=FONT_DISPLAY)


# ═══════════════════════════════════════════════════════
#  SLIDE TITRE – Design splitté
# ═══════════════════════════════════════════════════════

def _make_title_slide(prs, specialite: str, module: str, chapitre: str, niveau: str,
                      title_image: bytes = None, photographer: str = ''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG)

    # ── Panneau gauche violet (gradient vertical) ─────
    panel = _rect(slide, 0, 0, 4.4, 7.5, C_PANEL)
    _fill_gradient(panel, [
        (0,   C_PANEL),
        (65,  C_PANEL),
        (100, C_ACCENT_DARK),
    ], angle_deg=90)
    # Ligne de séparation subtile au raccord
    _rect(slide, 0, 5.58, 4.4, 0.04, C_ACCENT2)

    # ── Déco panneau gauche ───────────────────────────
    circ = slide.shapes.add_shape(9, Inches(2.8), Inches(-1.0), Inches(3.0), Inches(3.0))
    circ.fill.solid()
    circ.fill.fore_color.rgb = C_ACCENT_DARK
    _no_line(circ)

    # ── Textes panneau gauche ─────────────────────────
    tb_spec = _tb(slide, 0.35, 1.4, 3.7, 0.5)
    tf = tb_spec.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = specialite.upper()
    _run_fmt(run, 13, C_WHITE, bold=True, font=FONT_DISPLAY)

    _rect(slide, 0.35, 2.0, 1.2, 0.04, C_WHITE)

    _rounded_rect(slide, 0.35, 2.2, 1.5, 0.38, C_ACCENT_DARK)
    tb_niv = _tb(slide, 0.45, 2.25, 1.3, 0.3)
    p2 = tb_niv.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = f"NIVEAU {niveau.upper()}"
    _run_fmt(run2, 11, C_ACCENT2, bold=True, font=FONT_DISPLAY)

    tb_mod = _tb(slide, 0.35, 5.75, 3.7, 0.9)
    tf_mod = tb_mod.text_frame
    tf_mod.word_wrap = True
    p3 = tf_mod.paragraphs[0]
    run3 = p3.add_run()
    run3.text = module
    _run_fmt(run3, 12, RGBColor(0xC7, 0xD2, 0xFE), italic=True)

    # ── Panneau droit : image ou déco ─────────────────
    RIGHT_L = 4.4
    RIGHT_W = 13.33 - RIGHT_L   # 8.93"

    if title_image:
        # Image de fond sur tout le panneau droit (pleine hauteur, pleine visibilité)
        slide.shapes.add_picture(
            BytesIO(title_image),
            Inches(RIGHT_L), Inches(0), Inches(RIGHT_W), Inches(7.5)
        )
        # Bande sombre UNIQUEMENT en bas (zone du titre) — l'image reste visible
        # sur toute la partie haute. Opacité élevée (92 %) pour que le titre blanc
        # soit parfaitement lisible quelle que soit l'image choisie.
        BAND_TOP = 5.0
        BAND_H   = 2.5
        _add_transparent_rect(slide, RIGHT_L, BAND_TOP, RIGHT_W, BAND_H, C_BG, 92)

        # Trait accent placé AU-DESSUS du titre, dans la bande sombre
        _rect(slide, 4.75, BAND_TOP + 0.22, 2.0, 0.05, C_ACCENT)

        # Titre principal centré verticalement dans la bande sombre
        tb_title = _tb(slide, 4.75, BAND_TOP + 0.45, 8.2, BAND_H - 0.75)
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        _add_runs(p_title, chapitre, 32, C_WHITE, base_bold=True, font=FONT_DISPLAY)
        p_title.alignment = PP_ALIGN.LEFT
        p_title.space_after = Pt(6)
    else:
        # Pas d'image : grand cercle déco + titre en haut (comportement historique)
        circ2 = slide.shapes.add_shape(9, Inches(9.5), Inches(4.5), Inches(4.5), Inches(4.5))
        circ2.fill.solid()
        circ2.fill.fore_color.rgb = C_ACCENT_DARK
        _no_line(circ2)

        tb_title = _tb(slide, 4.75, 1.5, 8.2, 4.2)
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        _add_runs(p_title, chapitre, 36, C_WHITE, base_bold=True, font=FONT_DISPLAY)
        p_title.alignment = PP_ALIGN.LEFT
        p_title.space_after = Pt(10)

        # Trait déco sous le titre
        _rect(slide, 4.75, 6.1, 2.0, 0.05, C_ACCENT)

    # Attribution photographe Unsplash (obligatoire)
    if title_image and photographer:
        tb_photo = _tb(slide, RIGHT_L + 0.15, 7.2, RIGHT_W - 0.3, 0.28)
        p_ph = tb_photo.text_frame.paragraphs[0]
        p_ph.alignment = PP_ALIGN.RIGHT
        run_ph = p_ph.add_run()
        run_ph.text = f"Photo: {photographer} / Unsplash"
        _run_fmt(run_ph, 7, C_TEXT_MUTED)

    return slide


# ═══════════════════════════════════════════════════════
#  SLIDE DE SECTION – Fond bicolore
# ═══════════════════════════════════════════════════════

def _make_section_slide(prs, title: str, numero: int = 0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_ACCENT)

    # Bande sombre en haut
    _rect(slide, 0, 0, 13.33, 1.1, C_ACCENT_DARK)
    # Bande sombre en bas
    _rect(slide, 0, 6.3, 13.33, 1.2, C_ACCENT_DARK)
    # Ligne de séparation haute
    _rect(slide, 0, 1.08, 13.33, 0.04, C_ACCENT2)
    # Ligne de séparation basse
    _rect(slide, 0, 6.28, 13.33, 0.04, C_ACCENT2)

    # Grand chiffre romain en fond (watermark)
    if 0 < numero <= len(CHIFFRES_ROMAINS):
        tb_num = _tb(slide, 7.5, 0.5, 5.5, 6.0)
        tf = tb_num.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = CHIFFRES_ROMAINS[numero - 1]
        _run_fmt(run, 200, RGBColor(0x7C, 0x7F, 0xF5), bold=True, font=FONT_DISPLAY)

    # Petite étiquette "SECTION X" en haut
    tb_lbl = _tb(slide, 0.6, 0.3, 5, 0.45)
    p_lbl = tb_lbl.text_frame.paragraphs[0]
    run_lbl = p_lbl.add_run()
    run_lbl.text = f"SECTION {CHIFFRES_ROMAINS[numero - 1]}" if 0 < numero <= len(CHIFFRES_ROMAINS) else "SECTION"
    _run_fmt(run_lbl, 11, RGBColor(0xC7, 0xD2, 0xFE), bold=True, font=FONT_DISPLAY)

    # Titre de section
    tb = _tb(slide, 0.6, 1.5, 10.0, 4.5)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_runs(p, title, 40, C_WHITE, base_bold=True, font=FONT_DISPLAY)
    p.alignment = PP_ALIGN.LEFT

    return slide


# ═══════════════════════════════════════════════════════
#  SLIDE SOMMAIRE (TOC cliquable)
# ═══════════════════════════════════════════════════════

def _make_toc_placeholder(prs):
    """
    Crée une slide vide « sommaire » juste après la slide titre.
    Sera remplie plus tard via _fill_toc_slide une fois que les slides de
    section ont été créées (pour pouvoir les référencer en hyperlien).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG_CARD)
    if BG_PATTERN_ENABLED:
        _add_pattern_overlay(slide)

    # Barres décoratives identiques aux slides de contenu
    _rect(slide, 0, 0, 0.09, 7.5, C_ACCENT)
    _rect(slide, 0.09, 0, 13.24, 0.05, C_ACCENT_MID)
    return slide


def _fill_toc_slide(slide, entries: list):
    """
    Remplit la slide sommaire : titre « Sommaire » + liste numérotée d'entrées
    cliquables. entries = [(titre_section, slide_cible_ou_None), ...].

    Mise en page adaptative :
      - ≤ 8 entrées : une colonne, taille de police élevée
      - 9-16 entrées : deux colonnes
      - > 16 : tronqué à 16
    """
    if not entries:
        return
    entries = entries[:16]
    n = len(entries)

    # ── Titre "Sommaire"
    tb_title = _tb(slide, 0.22, 0.35, 12.9, 1.0)
    tf_t = tb_title.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    _add_runs(p_t, "Sommaire", 34, C_WHITE, base_bold=True, font=FONT_DISPLAY)

    # Sous-titre discret
    tb_sub = _tb(slide, 0.22, 1.10, 12.9, 0.35)
    p_sub = tb_sub.text_frame.paragraphs[0]
    run_sub = p_sub.add_run()
    run_sub.text = f"{n} section{'s' if n > 1 else ''} · cliquez pour y accéder"
    _run_fmt(run_sub, 11, C_TEXT_MUTED, font=FONT_BODY)

    # Séparateur
    _rect(slide, 0.22, 1.55, 2.8, 0.05, C_ACCENT)

    body_top = 1.95
    body_h   = 5.10

    if n <= 8:
        # Une colonne centrale
        col_x  = 0.9
        col_w  = 11.5
        entry_h = body_h / n
        # Police adaptative : 22pt pour 3 entrées → 15pt pour 8
        font_size = max(14, int(22 - (n - 3) * 1.2)) if n >= 3 else 22
        for i, (titre, target) in enumerate(entries):
            y = body_top + i * entry_h
            _draw_toc_entry(slide, col_x, y, col_w, entry_h, i + 1, titre,
                            target, font_size=font_size)
    else:
        # Deux colonnes
        per_col = (n + 1) // 2
        col_w   = 6.1
        gap     = 0.35
        start_x = (13.33 - (2 * col_w + gap)) / 2
        entry_h = body_h / per_col
        font_size = 13
        for i, (titre, target) in enumerate(entries):
            col = i // per_col
            row = i % per_col
            x = start_x + col * (col_w + gap)
            y = body_top + row * entry_h
            _draw_toc_entry(slide, x, y, col_w, entry_h, i + 1, titre,
                            target, font_size=font_size)


def _draw_toc_entry(slide, x: float, y: float, w: float, h: float,
                    numero: int, titre: str, target_slide,
                    font_size: int = 18):
    """
    Dessine une ligne de sommaire : chip numéro romain + titre cliquable.
    target_slide peut être None (entrée non-cliquable, fallback gracieux).
    """
    # Chip avec le chiffre romain
    chip_diam = min(0.55, h * 0.72)
    chip_y    = y + (h - chip_diam) / 2
    roman     = (CHIFFRES_ROMAINS[numero - 1]
                 if 1 <= numero <= len(CHIFFRES_ROMAINS) else str(numero))
    _icon_chip(slide, x, chip_y, chip_diam, roman, bg_color=C_ACCENT,
               glyph_size_pt=max(11, int(chip_diam * 22)))

    # Titre à droite du chip
    text_left = x + chip_diam + 0.22
    text_w    = w - (chip_diam + 0.22)
    tb = _tb(slide, text_left, y, text_w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.04)
    tf.margin_right  = Inches(0.04)
    tf.margin_top    = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = _truncate(titre, 80)
    _run_fmt(run, font_size, C_TEXT_LIGHT, bold=True, font=FONT_DISPLAY)

    if target_slide is not None:
        _add_slide_hyperlink(run, slide, target_slide)
        # Indicateur visuel discret « → » à droite pour signaler la cliquabilité
        tb_arrow = _tb(slide, x + w - 0.45, y, 0.4, h)
        tf_a = tb_arrow.text_frame
        tf_a.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_a.margin_left = tf_a.margin_right = Inches(0)
        p_a = tf_a.paragraphs[0]
        p_a.alignment = PP_ALIGN.RIGHT
        run_a = p_a.add_run()
        run_a.text = '→'
        _run_fmt(run_a, max(12, font_size - 2), C_ACCENT2, bold=False,
                 font=FONT_DISPLAY)
        # L'icône elle aussi pointe vers la même cible (plus cliquable)
        _add_slide_hyperlink(run_a, slide, target_slide)


def _move_slide_to(prs, slide, new_index: int):
    """
    Déplace une slide dans le sldIdLst pour changer son ordre d'apparition.
    Utilisé pour positionner la slide sommaire juste après la slide titre.
    """
    sldIdLst = prs.slides._sldIdLst
    # Récupère le <p:sldId> correspondant à notre slide via r:id
    target_rId = None
    for rel in prs.part.rels.values():
        if rel.target_part is slide.part:
            target_rId = rel.rId
            break
    if target_rId is None:
        return
    target_sldId = None
    for sldId in list(sldIdLst):
        if sldId.get(qn('r:id')) == target_rId:
            target_sldId = sldId
            break
    if target_sldId is None:
        return
    sldIdLst.remove(target_sldId)
    sldIdLst.insert(new_index, target_sldId)


def _remove_slide(prs, slide):
    """
    Supprime une slide du deck : retire l'entrée <p:sldId> du sldIdLst et
    détache la relation vers la part du slide. La part elle-même reste dans
    le package mais n'est plus référencée (sera ignorée par PowerPoint).
    """
    sldIdLst = prs.slides._sldIdLst
    target_rId = None
    for rel in list(prs.part.rels.values()):
        if rel.target_part is slide.part:
            target_rId = rel.rId
            break
    if target_rId is None:
        return
    for sldId in list(sldIdLst):
        if sldId.get(qn('r:id')) == target_rId:
            sldIdLst.remove(sldId)
            break
    try:
        prs.part.drop_rel(target_rId)
    except KeyError:
        pass


# ═══════════════════════════════════════════════════════
#  SLIDE DE CONTENU – Barre gauche + badge section
# ═══════════════════════════════════════════════════════

def _content_base(prs, title: str, section_label: str = ''):
    """Crée la structure de base d'une slide de contenu (fond + déco)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG_CARD)

    # Motif de fond subtil (première forme = sous tout le reste dans la z-order)
    if BG_PATTERN_ENABLED:
        _add_pattern_overlay(slide)

    # Barre gauche accent
    _rect(slide, 0, 0, 0.09, 7.5, C_ACCENT)

    # Barre horizontale haute (fine)
    _rect(slide, 0.09, 0, 13.24, 0.05, C_ACCENT_MID)

    # Badge section
    if section_label:
        _section_badge(slide, section_label, 0.22, 0.12)

    return slide


def _make_content_slide(prs, title: str, content_lines: list[str],
                        is_bullets: bool = True, section_label: str = ''):
    slide = _content_base(prs, title, section_label)

    # Titre — autofit pour gérer les chapitres au libellé long
    top_title = 0.5 if section_label else 0.18
    tb_title = _tb(slide, 0.22, top_title, 12.9, 1.0)
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p = tf_title.paragraphs[0]
    _add_runs(p, title, 24, C_WHITE, base_bold=True, font=FONT_DISPLAY)
    _enable_shrink_to_fit(tf_title)

    # Séparateur
    sep_top = top_title + 1.0
    _rect(slide, 0.22, sep_top, 12.9, 0.025, C_ACCENT)

    # Corps — overflow protection : si la liste de bullets / le paragraphe
    # déborde, PowerPoint réduira automatiquement la police au rendu.
    body_top = sep_top + 0.18
    body_h   = 7.5 - body_top - 0.35
    tb_body = _tb(slide, 0.22, body_top, 12.9, body_h)
    tf = tb_body.text_frame
    tf.word_wrap = True

    for i, line in enumerate(content_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5 if is_bullets else 3)
        p.space_after  = Pt(2)
        if is_bullets:
            run_b = p.add_run()
            run_b.text = '◆  '
            _run_fmt(run_b, 9, C_BULLET, bold=True)
            _add_runs(p, line, 13, C_TEXT_LIGHT)
        else:
            _add_runs(p, line, 13, C_TEXT_LIGHT)

    _enable_shrink_to_fit(tf)
    return slide


def _make_two_column_slide(prs, title: str, bullets: list[str], section_label: str = ''):
    slide = _content_base(prs, title, section_label)

    top_title = 0.5 if section_label else 0.18
    _icon_chip(slide, 0.22, top_title + 0.12, 0.5, CHIP_COMPARE, bg_color=C_ACCENT)
    tb_title = _tb(slide, 0.85, top_title, 12.27, 1.0)
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p = tf_title.paragraphs[0]
    _add_runs(p, title, 24, C_WHITE, base_bold=True, font=FONT_DISPLAY)
    _enable_shrink_to_fit(tf_title)

    sep_top = top_title + 1.0
    _rect(slide, 0.22, sep_top, 12.9, 0.025, C_ACCENT)

    body_top = sep_top + 0.18
    body_h   = 7.5 - body_top - 0.35

    mid = (len(bullets) + 1) // 2
    col1, col2 = bullets[:mid], bullets[mid:]

    # Séparateur vertical central
    _rect(slide, 6.8, body_top, 0.02, body_h, C_ACCENT_MID)

    # Chaque colonne : autofit indépendant — si une colonne déborde plus
    # que l'autre (cas asymétrique), seule celle-ci est réduite.
    for col_lines, lx in [(col1, 0.22), (col2, 6.95)]:
        tb = _tb(slide, lx, body_top, 6.35, body_h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(col_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(5)
            p.space_after  = Pt(2)
            run_b = p.add_run()
            run_b.text = '◆  '
            _run_fmt(run_b, 9, C_BULLET, bold=True)
            _add_runs(p, line, 13, C_TEXT_LIGHT)
        _enable_shrink_to_fit(tf)

    return slide


# ═══════════════════════════════════════════════════════
#  SLIDE TABLEAU
# ═══════════════════════════════════════════════════════

_NUMERIC_CELL_RE = re.compile(
    r'^\s*[-+]?\d[\d\s.,]*\s*(%|€|\$|£|¥|pts?|x|×|k|M|Md|Mds|h|j|an|ans|°C|°|\/\d+)?\s*$',
    re.IGNORECASE,
)


def _is_numeric_cell(text: str) -> bool:
    """True si la cellule ressemble à une valeur numérique (avec suffixe éventuel)."""
    if not text or not text.strip():
        return False
    return bool(_NUMERIC_CELL_RE.match(text.strip()))


def _remove_cell_borders(cell):
    """
    Supprime les bordures par défaut des cellules de tableau.
    python-pptx ne fournit pas d'API pour ça — on passe par l'XML :
    on ajoute <a:ln><a:noFill/></a:ln> sur chaque arête (L/R/T/B).
    Rendu final : séparateurs visuels assurés uniquement par les couleurs
    de remplissage alternées, plus propre et plus moderne.
    """
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ('lnL', 'lnR', 'lnT', 'lnB'):
        # Retire les éventuels <a:lnX> existants avant d'ajouter les nôtres
        for existing in tcPr.findall(qn(f'a:{tag}')):
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(f'a:{tag}'))
        ln.set('w', '0')
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        etree.SubElement(ln, qn('a:noFill'))


def _make_table_slide(prs, title: str, headers: list[str], rows: list[list[str]],
                      section_label: str = ''):
    slide = _content_base(prs, title, section_label)

    top_title = 0.5 if section_label else 0.18
    _icon_chip(slide, 0.22, top_title + 0.12, 0.5, CHIP_TABLE, bg_color=C_ACCENT)
    tb_title = _tb(slide, 0.85, top_title, 12.27, 1.0)
    tf = tb_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_runs(p, title, 24, C_WHITE, base_bold=True, font=FONT_DISPLAY)
    _enable_shrink_to_fit(tf)

    sep_top = top_title + 1.0
    _rect(slide, 0.22, sep_top, 12.9, 0.025, C_ACCENT)

    n_cols = max(len(headers), max((len(r) for r in rows), default=0))
    n_rows = min(len(rows) + 1, 12)
    if n_cols == 0 or n_rows < 2:
        return slide

    # Détection des colonnes numériques : si >=60% des valeurs d'une colonne
    # ressemblent à des nombres, on aligne à droite toute la colonne.
    numeric_cols = set()
    for j in range(n_cols):
        col_vals = [(rows[i][j] if j < len(rows[i]) else '') for i in range(len(rows))]
        non_empty = [v for v in col_vals if v and v.strip()]
        if non_empty and sum(1 for v in non_empty if _is_numeric_cell(v)) >= max(2, int(0.6 * len(non_empty))):
            numeric_cols.add(j)

    # Première colonne en gras si c'est typiquement une colonne label (texte court,
    # pas numérique). Évite d'ajouter du gras partout.
    bold_first_col = (0 not in numeric_cols) and all(
        len((rows[i][0] if 0 < len(rows[i]) else '')) < 40 for i in range(len(rows))
    )

    tbl_top  = Inches(sep_top + 0.22)
    tbl_left = Inches(0.22)
    tbl_w    = Inches(12.9)
    # Hauteur header légèrement plus grande que les rows pour mieux marquer la hiérarchie
    header_h = Inches(min(0.58, 6.0 / n_rows + 0.05))
    body_row_h = Inches(min(0.50, 5.6 / max(n_rows - 1, 1)))
    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, tbl_left, tbl_top, tbl_w,
        header_h + body_row_h * (n_rows - 1),
    )
    tbl = tbl_shape.table

    # Hauteurs de ligne explicites (header distinct + body réguliers)
    tbl.rows[0].height = header_h
    for i in range(1, n_rows):
        tbl.rows[i].height = body_row_h

    col_w = tbl_w // n_cols
    for col in tbl.columns:
        col.width = col_w

    def _cell(cell, bg: RGBColor, text: str, *, bold=False, align=PP_ALIGN.LEFT,
              color=C_WHITE, size=11, is_header=False):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        # Padding cellule : aère un peu le contenu
        cell.margin_left   = Inches(0.10)
        cell.margin_right  = Inches(0.10)
        cell.margin_top    = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        _remove_cell_borders(cell)

        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after  = Pt(0)
        _add_runs(p, str(text), size, color, base_bold=bold,
                  font=(FONT_DISPLAY if is_header else FONT_BODY))
        # Cellules de tableau : si le texte d'une cellule déborde de la
        # hauteur de ligne, PowerPoint réduit la police au rendu.
        _enable_shrink_to_fit(tf)

    # Header row : alignement selon colonne (numérique → droite, sinon gauche)
    for j, h in enumerate(headers[:n_cols]):
        align = PP_ALIGN.RIGHT if j in numeric_cols else PP_ALIGN.LEFT
        # Centre la 1re colonne label si elle est un index court (<= 3 lettres)
        if j == 0 and bold_first_col and len(str(h).strip()) <= 3:
            align = PP_ALIGN.LEFT
        _cell(tbl.cell(0, j), C_TABLE_HDR, h, bold=True, align=align,
              color=C_WHITE, size=12, is_header=True)

    # Body rows : alternance + alignement numérique à droite
    for i, row in enumerate(rows[:n_rows - 1]):
        bg = C_TABLE_ROW1 if i % 2 == 0 else C_TABLE_ROW2
        for j in range(n_cols):
            val = row[j] if j < len(row) else ''
            is_num = j in numeric_cols
            align = PP_ALIGN.RIGHT if is_num else PP_ALIGN.LEFT
            bold  = bold_first_col and j == 0
            color = C_ACCENT2 if (bold and j == 0) else C_TEXT_LIGHT
            _cell(tbl.cell(i + 1, j), bg, val, bold=bold, align=align,
                  color=color, size=11)

    return slide


# ═══════════════════════════════════════════════════════
#  SLIDE DÉFINITIONS
# ═══════════════════════════════════════════════════════

def _make_definitions_slide(prs, items: list[tuple[str, str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG_CARD)
    _rect(slide, 0, 0, 0.09, 7.5, C_ACCENT)

    # Header band (gradient horizontal)
    header = _rect(slide, 0.09, 0, 13.24, 0.75, C_ACCENT_DARK)
    _fill_gradient(header, [(0, C_ACCENT_DARK), (100, C_ACCENT_MID)], angle_deg=0)

    _icon_chip(slide, 0.28, 0.14, 0.48, CHIP_DEF, bg_color=C_ACCENT)
    tb_hdr = _tb(slide, 0.9, 0.1, 12.1, 0.55)
    p = tb_hdr.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = 'DÉFINITIONS DES CONCEPTS CLÉS'
    _run_fmt(run, 16, C_WHITE, bold=True, font=FONT_DISPLAY)

    # Ligne sous le header
    _rect(slide, 0.09, 0.73, 13.24, 0.03, C_ACCENT)

    tb_body = _tb(slide, 0.3, 0.9, 12.85, 6.2)
    tf = tb_body.text_frame
    tf.word_wrap = True

    for i, (terme, definition) in enumerate(items[:10]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        # Puce + terme
        run_b = p.add_run()
        run_b.text = '◆ '
        _run_fmt(run_b, 10, C_BULLET, bold=True)
        run_t = p.add_run()
        run_t.text = f"{_clean(terme)}  "
        _run_fmt(run_t, 13, C_ACCENT2, bold=True, font=FONT_DISPLAY)
        run_d = p.add_run()
        run_d.text = _clean(definition)
        _run_fmt(run_d, 12, C_TEXT_LIGHT)

    _enable_shrink_to_fit(tf)
    return slide


# ═══════════════════════════════════════════════════════
#  SLIDE POINTS IMPORTANTS – Numéros cerclés
# ═══════════════════════════════════════════════════════

def _make_key_points_slide(prs, points: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG_CARD)
    _rect(slide, 0, 0, 0.09, 7.5, C_ACCENT)

    # Header band (gradient horizontal)
    header = _rect(slide, 0.09, 0, 13.24, 0.75, C_ACCENT_DARK)
    _fill_gradient(header, [(0, C_ACCENT_DARK), (100, C_ACCENT_MID)], angle_deg=0)

    _icon_chip(slide, 0.28, 0.14, 0.48, CHIP_PTS, bg_color=C_ACCENT)
    tb_hdr = _tb(slide, 0.9, 0.1, 12.1, 0.55)
    p = tb_hdr.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'POINTS IMPORTANTS À RETENIR'
    _run_fmt(run, 16, C_WHITE, bold=True, font=FONT_DISPLAY)

    _rect(slide, 0.09, 0.73, 13.24, 0.03, C_ACCENT)

    # Mise en page 2 colonnes si > 6 points
    pts = points[:12]
    if len(pts) > 6:
        mid = (len(pts) + 1) // 2
        cols = [(pts[:mid], 0.3), (pts[mid:], 6.85)]
        _rect(slide, 6.76, 0.85, 0.02, 6.3, C_ACCENT_MID)
    else:
        cols = [(pts, 0.3)]

    for col_pts, lx in cols:
        col_w = 6.2 if len(cols) > 1 else 12.85
        tb = _tb(slide, lx, 0.9, col_w, 6.2)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, point in enumerate(col_pts):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(6)
            p.space_after  = Pt(2)
            idx = pts.index(point)
            run_num = p.add_run()
            run_num.text = f"{CIRCLED_NUMS[idx]}  "
            _run_fmt(run_num, 15, C_ACCENT, bold=True, font=FONT_DISPLAY)
            _add_runs(p, point, 13, C_TEXT_LIGHT)
        _enable_shrink_to_fit(tf)

    return slide


def _make_callout_slide(prs, quote: str, attribution: str = '',
                         section_label: str = ''):
    """
    Slide "callout / citation" : un texte court mis en exergue (citation,
    règle clé, principe fondamental). Pensée pour créer une rupture visuelle
    et faire pause dans la lecture.

    Rendu : fond dégradé violet sombre + grand guillemet décoratif en haut
    gauche + citation en italique centrée + attribution en bas.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG)

    # ── Fond dégradé pleine slide (diagonale violet sombre → noir) ──
    bg_rect = _rect(slide, 0, 0, 13.33, 7.5, C_ACCENT_DARK)
    _fill_gradient(bg_rect, [
        (0,   C_ACCENT_DARK),
        (60,  RGBColor(0x12, 0x13, 0x38)),
        (100, C_BG),
    ], angle_deg=135)

    # ── Barre accent gauche (cohérence avec autres slides) ──
    _rect(slide, 0, 0, 0.09, 7.5, C_ACCENT)
    _rect(slide, 0.09, 0, 13.24, 0.05, C_ACCENT_MID)

    # Badge section (optionnel)
    if section_label:
        _section_badge(slide, section_label, 0.22, 0.12)

    # ── Grand guillemet décoratif (coin haut gauche) ──
    tb_quote_deco = _tb(slide, 0.5, 0.55, 2.0, 2.0)
    tf_qd = tb_quote_deco.text_frame
    tf_qd.margin_left = 0
    tf_qd.margin_top = 0
    p_qd = tf_qd.paragraphs[0]
    run_qd = p_qd.add_run()
    run_qd.text = '«'
    _run_fmt(run_qd, 200, C_ACCENT_MID, bold=True, font=FONT_DISPLAY)

    # ── Taille adaptative de la citation selon longueur ──
    n = len(quote)
    if n < 80:
        quote_pt = 34
    elif n < 150:
        quote_pt = 28
    elif n < 250:
        quote_pt = 22
    else:
        quote_pt = 18
        quote = _truncate(quote, 400)

    # ── Citation centrée ──
    # Zone large, avec marges gauche/droite généreuses pour respirer
    tb_quote = _tb(slide, 1.6, 2.0, 10.1, 3.5)
    tf_q = tb_quote.text_frame
    tf_q.word_wrap = True
    tf_q.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_q = tf_q.paragraphs[0]
    p_q.alignment = PP_ALIGN.CENTER
    run_q = p_q.add_run()
    run_q.text = quote
    _run_fmt(run_q, quote_pt, C_WHITE, italic=True, font=FONT_DISPLAY)
    # La citation a déjà une taille adaptative selon longueur ; autofit en
    # complément pour les cas extrêmes (citations très longues + saut de ligne).
    _enable_shrink_to_fit(tf_q)

    # ── Trait accent sous la citation ──
    _rect(slide, 6.16, 5.75, 1.0, 0.04, C_ACCENT)

    # ── Attribution en bas (italique, muted, préfixée par —) ──
    if attribution:
        tb_attr = _tb(slide, 1.6, 5.9, 10.1, 0.6)
        tf_a = tb_attr.text_frame
        tf_a.word_wrap = True
        p_a = tf_a.paragraphs[0]
        p_a.alignment = PP_ALIGN.CENTER
        run_a = p_a.add_run()
        # « — » (em dash) : typographie française propre pour attribution
        run_a.text = f'— {attribution}'
        _run_fmt(run_a, 14, C_TEXT_MUTED, italic=True, font=FONT_BODY)

    return slide


def _split_objective_phrase(text: str) -> tuple[str, str]:
    """
    Sépare un objectif pédagogique en (verbe_gras, reste).
    Si le texte commence par **verbe** (ex : "**Analyser** les flux…"), on met
    le verbe en couleur accent pour souligner la compétence visée.
    Sinon retourne ('', text).
    """
    m = re.match(r'^\*\*(.+?)\*\*\s*(.*)', text.strip())
    if m:
        return _clean(m.group(1)), m.group(2).strip()
    return '', text.strip()


def _make_objectives_slide(prs, items: list[str], section_label: str = ''):
    """
    Slide « Objectifs pédagogiques ».

    Placée en tête de chapitre ou de section, elle annonce les compétences
    visées. Rendu : bandeau titre « OBJECTIFS PÉDAGOGIQUES » + phrase
    d'introduction « À l'issue de cette séquence, vous serez capable de : »
    + liste de 3 à 6 objectifs, chacun précédé d'une pastille ✓ verte.

    Les items de la forme « **Verbe** complément » voient leur verbe mis en
    couleur accent (souligne le verbe d'action pédagogique).
    """
    # Normalisation
    objs = [_clean(str(it)) for it in (items or []) if str(it).strip()]
    if not objs:
        return None
    objs = objs[:6]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG_CARD)
    if BG_PATTERN_ENABLED:
        _add_pattern_overlay(slide)
    _rect(slide, 0, 0, 0.09, 7.5, C_ACCENT)

    # ── En-tête : bande dégradée + chip cible + label
    header = _rect(slide, 0.09, 0, 13.24, 0.85, C_ACCENT_DARK)
    _fill_gradient(header, [(0, C_ACCENT_DARK), (100, C_ACCENT_MID)], angle_deg=0)
    _icon_chip(slide, 0.28, 0.19, 0.48, CHIP_GOAL, bg_color=C_ACCENT)
    tb_hdr = _tb(slide, 0.9, 0.15, 12.1, 0.55)
    p_h = tb_hdr.text_frame.paragraphs[0]
    run_h = p_h.add_run()
    run_h.text = 'OBJECTIFS PÉDAGOGIQUES'
    _run_fmt(run_h, 16, C_WHITE, bold=True, font=FONT_DISPLAY)
    _rect(slide, 0.09, 0.83, 13.24, 0.03, C_SUCCESS)

    # ── Phrase d'introduction
    tb_intro = _tb(slide, 0.5, 1.0, 12.33, 0.55)
    tf_i = tb_intro.text_frame
    tf_i.word_wrap = True
    p_i = tf_i.paragraphs[0]
    run_i = p_i.add_run()
    run_i.text = 'À l’issue de cette séquence, vous serez capable de :'
    _run_fmt(run_i, 14, C_TEXT_MUTED, italic=True, font=FONT_BODY)

    # ── Dimensions adaptatives selon le nombre d'objectifs
    n = len(objs)
    body_top = 1.70
    body_bot = 7.15
    body_h   = body_bot - body_top
    row_h    = body_h / n

    # Chip/police proportionnels (3 → large et aéré ; 6 → compact)
    if n <= 3:
        chip_d, chip_pt, txt_pt = 0.60, 20, 18
    elif n == 4:
        chip_d, chip_pt, txt_pt = 0.52, 18, 16
    elif n == 5:
        chip_d, chip_pt, txt_pt = 0.46, 16, 15
    else:  # 6
        chip_d, chip_pt, txt_pt = 0.40, 14, 14

    chip_left = 0.55
    text_left = chip_left + chip_d + 0.30
    text_w    = 13.33 - text_left - 0.40

    for i, obj in enumerate(objs):
        y_center = body_top + row_h * (i + 0.5)
        y_chip   = y_center - chip_d / 2

        # Pastille verte avec coche
        _icon_chip(slide, chip_left, y_chip, chip_d, CHIP_CHECK,
                   bg_color=C_SUCCESS, glyph_size_pt=chip_pt)

        # Texte de l'objectif — verbe d'action mis en couleur accent si détecté
        tb = _tb(slide, text_left, y_chip - 0.05, text_w, chip_d + 0.10)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left   = Inches(0.04)
        tf.margin_right  = Inches(0.04)
        tf.margin_top    = Inches(0)
        tf.margin_bottom = Inches(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]

        verb, rest = _split_objective_phrase(obj)
        if verb:
            rv = p.add_run()
            rv.text = verb + ' '
            _run_fmt(rv, txt_pt, C_ACCENT2, bold=True, font=FONT_DISPLAY)
            _add_runs(p, rest, txt_pt, C_TEXT_LIGHT)
        else:
            _add_runs(p, obj, txt_pt, C_TEXT_LIGHT)
        _enable_shrink_to_fit(tf)

    return slide


def _make_synthese_slide(prs, items: list[str], section_label: str = ''):
    """
    Slide « Synthèse / À retenir ».

    Placée en fin de chapitre ou de section, elle récapitule les 3 à 5
    points clés. Rendu : bandeau « SYNTHÈSE • À RETENIR » + cartes
    numérotées avec chip étoile, titre en gras et corps en texte clair.

    Les items de la forme « **Lead** : explication » sont rendus en deux
    temps : lead en couleur accent + corps sous forme de phrase support.
    """
    pts = [_clean(str(it)) for it in (items or []) if str(it).strip()]
    if not pts:
        return None
    pts = pts[:5]
    n = len(pts)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C_BG_CARD)
    if BG_PATTERN_ENABLED:
        _add_pattern_overlay(slide)
    _rect(slide, 0, 0, 0.09, 7.5, C_ACCENT)

    # ── En-tête dégradé avec chip synthèse
    header = _rect(slide, 0.09, 0, 13.24, 0.85, C_ACCENT_DARK)
    _fill_gradient(header, [(0, C_ACCENT_DARK), (100, C_ACCENT_MID)], angle_deg=0)
    _icon_chip(slide, 0.28, 0.19, 0.48, CHIP_RECAP, bg_color=C_ACCENT)
    tb_hdr = _tb(slide, 0.9, 0.15, 12.1, 0.55)
    p_h = tb_hdr.text_frame.paragraphs[0]
    run_h = p_h.add_run()
    run_h.text = 'SYNTHÈSE  •  À RETENIR'
    _run_fmt(run_h, 16, C_WHITE, bold=True, font=FONT_DISPLAY)
    _rect(slide, 0.09, 0.83, 13.24, 0.03, C_ACCENT)

    # ── Zone cartes (layout vertical, une carte par take-away)
    body_top = 1.15
    body_bot = 7.15
    body_h   = body_bot - body_top
    gap      = 0.15
    card_h   = (body_h - gap * (n - 1)) / n

    # Police adaptative : 3 → confortable ; 5 → compacte
    if n <= 3:
        num_pt, lead_pt, body_pt = 32, 18, 14
        lead_max = 300
    elif n == 4:
        num_pt, lead_pt, body_pt = 26, 16, 13
        lead_max = 200
    else:  # 5
        num_pt, lead_pt, body_pt = 22, 15, 12
        lead_max = 150

    card_left = 0.40
    card_w    = 13.33 - card_left - 0.40
    num_w     = 0.90  # largeur de la colonne numéro

    for i, point in enumerate(pts):
        y = body_top + i * (card_h + gap)

        # Carte fond dégradé (légère surcouche visuelle)
        card = _rounded_rect(slide, card_left, y, card_w, card_h, C_ACCENT_DARK)
        _fill_gradient(card, [(0, C_ACCENT_DARK), (100, C_BG_CARD)], angle_deg=0)
        _add_shadow(card, blur_pt=8, dist_pt=2, alpha_pct=35, dir_deg=90)

        # Liséré vertical accent sur le bord gauche de la carte
        _rect(slide, card_left, y, 0.06, card_h, C_ACCENT)

        # Numéro gros caractère dans la colonne de gauche
        tb_num = _tb(slide, card_left + 0.12, y, num_w, card_h)
        tf_n = tb_num.text_frame
        tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_n.margin_left = Inches(0.04)
        tf_n.margin_top  = Inches(0)
        p_n = tf_n.paragraphs[0]
        p_n.alignment = PP_ALIGN.CENTER
        run_n = p_n.add_run()
        run_n.text = str(i + 1).zfill(2)
        _run_fmt(run_n, num_pt, C_ACCENT2, bold=True, font=FONT_DISPLAY)

        # Texte principal (lead éventuel + corps)
        text_left = card_left + 0.20 + num_w + 0.15
        text_w    = card_w - (text_left - card_left) - 0.25
        tb_t = _tb(slide, text_left, y + 0.10, text_w, card_h - 0.20)
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left   = Inches(0.04)
        tf_t.margin_right  = Inches(0.04)
        tf_t.margin_top    = Inches(0)
        tf_t.margin_bottom = Inches(0)
        tf_t.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Détecte un lead « **Titre** : corps » ou « Titre : corps »
        lead, rest = '', _truncate(point, lead_max)
        mb = re.match(r'^\*\*(.+?)\*\*\s*[:：]\s*(.+)', point)
        if mb:
            lead = _clean(mb.group(1))
            rest = _truncate(_clean(mb.group(2)), lead_max)
        elif ':' in point and len(point.split(':', 1)[0]) <= 60:
            tit, body = point.split(':', 1)
            if body.strip():
                lead = _clean(tit)
                rest = _truncate(_clean(body), lead_max)

        if lead:
            p_l = tf_t.paragraphs[0]
            p_l.space_after = Pt(3)
            _add_runs(p_l, lead, lead_pt, C_WHITE, base_bold=True, font=FONT_DISPLAY)
            p_b = tf_t.add_paragraph()
            p_b.space_before = Pt(2)
            _add_runs(p_b, rest, body_pt, C_TEXT_LIGHT)
        else:
            p_l = tf_t.paragraphs[0]
            _add_runs(p_l, rest, lead_pt, C_WHITE, base_bold=True, font=FONT_DISPLAY)
        _enable_shrink_to_fit(tf_t)

    return slide


def _make_code_slide(prs, title: str, code: str, language: str = '',
                       description: str = '', section_label: str = ''):
    """
    Slide « Bloc de code » : encadré sombre avec liséré accent à gauche,
    police Consolas, étiquette de langage en haut à droite, optionnellement
    une phrase d'introduction au-dessus du code.

    Mise en page :
      ┌─[python]──────────────────────────┐
      │ ▌ def calculer_taux(brut):        │
      │ ▌     return brut * 0.85          │
      │ ▌                                 │
      │ ▌ # exemple                       │  ← lignes commentaires en muted
      │ ▌ net = calculer_taux(2500)       │
      └───────────────────────────────────┘

    Taille de police adaptative :
      ≤12 lignes & ≤70 chars/ligne → 14pt
      ≤20 lignes & ≤90 chars       → 12pt
      ≤30 lignes & ≤110 chars      → 10pt
      sinon                         → 9pt + truncate
    """
    code = code or ''
    if not code.strip():
        return None

    slide = _content_base(prs, title, section_label)

    # ── En-tête : chip clavier + titre + séparateur ──
    top_title = 0.5 if section_label else 0.18
    _icon_chip(slide, 0.22, top_title + 0.12, 0.5, CHIP_CODE,
               bg_color=C_ACCENT)
    tb_title = _tb(slide, 0.85, top_title, 12.27, 1.0)
    tf = tb_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_runs(p, title, 24, C_WHITE, base_bold=True, font=FONT_DISPLAY)
    _enable_shrink_to_fit(tf)

    sep_top = top_title + 1.0
    _rect(slide, 0.22, sep_top, 12.9, 0.025, C_ACCENT)

    body_top = sep_top + 0.30

    # ── Description (intro optionnelle au-dessus du code) ──
    if description:
        desc_h = 0.65
        tb_desc = _tb(slide, 0.30, body_top, 12.73, desc_h)
        tf_d = tb_desc.text_frame
        tf_d.word_wrap = True
        p_d = tf_d.paragraphs[0]
        _add_runs(p_d, _truncate(description, 200), 12, C_TEXT_MUTED)
        _enable_shrink_to_fit(tf_d)
        body_top += desc_h + 0.10

    # ── Encadré code ──
    code_top   = body_top
    code_bot   = 7.10
    code_left  = 0.30
    code_w     = 13.33 - 2 * code_left
    code_h     = code_bot - code_top

    # Fond légèrement plus sombre que C_BG_CARD pour distinguer le code
    code_bg_color = RGBColor(0x10, 0x12, 0x22)
    code_box = _rounded_rect(slide, code_left, code_top, code_w, code_h,
                              code_bg_color)
    _add_shadow(code_box, blur_pt=10, dist_pt=3, alpha_pct=45, dir_deg=90)

    # Liséré gauche accent (signature visuelle « bloc de code »)
    bar_w = 0.08
    _rect(slide, code_left, code_top, bar_w, code_h, C_ACCENT)

    # ── Étiquette langage (top-right pill) ──
    if language:
        lbl = language.lower()
        lw = max(0.65, len(lbl) * 0.10 + 0.20)
        lh = 0.30
        lx = code_left + code_w - lw - 0.15
        ly = code_top + 0.15
        pill = _rounded_rect(slide, lx, ly, lw, lh, C_ACCENT_DARK)
        tb_lang = _tb(slide, lx, ly, lw, lh)
        tf_l = tb_lang.text_frame
        tf_l.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_l.margin_left  = Inches(0.04)
        tf_l.margin_right = Inches(0.04)
        tf_l.margin_top    = Inches(0)
        tf_l.margin_bottom = Inches(0)
        p_l = tf_l.paragraphs[0]
        p_l.alignment = PP_ALIGN.CENTER
        run_l = p_l.add_run()
        run_l.text = lbl.upper()
        _run_fmt(run_l, 8, C_ACCENT2, bold=True, font=FONT_DISPLAY)

    # ── Texte du code ──
    text_left_pad = bar_w + 0.20
    text_top_pad  = 0.30 if language else 0.20
    code_text_x   = code_left + text_left_pad
    code_text_y   = code_top + text_top_pad
    code_text_w   = code_w - text_left_pad - 0.20
    code_text_h   = code_h - text_top_pad - 0.20

    # Adaptation taille police selon lignes/longueur
    lines = code.split('\n')
    n_lines = len(lines)
    max_line_len = max((len(l) for l in lines), default=0)

    if n_lines <= 12 and max_line_len <= 70:
        code_pt = 14
    elif n_lines <= 20 and max_line_len <= 90:
        code_pt = 12
    elif n_lines <= 30 and max_line_len <= 110:
        code_pt = 10
    else:
        code_pt = 9

    # Cap à 35 lignes — au-delà, ajout d'un marqueur de troncature
    truncated = False
    if n_lines > 35:
        lines = lines[:34] + ['... (suite tronquée)']
        truncated = True

    tb_code = _tb(slide, code_text_x, code_text_y, code_text_w, code_text_h)
    tf_c = tb_code.text_frame
    # On NE wrap PAS : si une ligne déborde, l'autofit la réduira via la
    # police globale. Wrap aurait casse l'indentation et la sémantique.
    tf_c.word_wrap = False
    tf_c.margin_left   = Inches(0.04)
    tf_c.margin_right  = Inches(0.04)
    tf_c.margin_top    = Inches(0.04)
    tf_c.margin_bottom = Inches(0.04)
    tf_c.vertical_anchor = MSO_ANCHOR.TOP

    for i, line in enumerate(lines):
        para = tf_c.paragraphs[0] if i == 0 else tf_c.add_paragraph()
        para.space_before = Pt(0)
        para.space_after  = Pt(0)
        run = para.add_run()
        # Ligne vide → espace pour préserver la hauteur visuelle
        run.text = line if line else ' '
        is_comment = _is_comment_line(line, language) if line.strip() else False
        if i == len(lines) - 1 and truncated:
            # Marqueur de troncature en italique muted
            _run_fmt(run, code_pt, C_TEXT_MUTED, italic=True, font=FONT_MONO)
        else:
            color = C_TEXT_MUTED if is_comment else C_TEXT_LIGHT
            _run_fmt(run, code_pt, color, font=FONT_MONO)

    _enable_shrink_to_fit(tf_c)
    return slide


def _make_case_study_slide(prs, title: str, contexte: str, probleme: str,
                             solution: str, resultat: str,
                             section_label: str = ''):
    """
    Slide « Cas pratique » : 4 cartes en grille 2x2.

      ┌─────────────────────┬─────────────────────┐
      │ 01 CONTEXTE   (mid) │ 02 PROBLÈME (warn)  │
      │ ...........        │ ...........        │
      ├─────────────────────┼─────────────────────┤
      │ 03 SOLUTION (acct)  │ 04 RÉSULTAT (succ)  │
      │ ...........        │ ...........        │
      └─────────────────────┴─────────────────────┘

    Lecture en Z (haut-gauche → haut-droit → bas-gauche → bas-droit) qui
    correspond à la narration d'un cas : on pose le décor, on identifie
    le problème, on raconte la solution, on mesure les résultats.

    Chaque carte a son bandeau supérieur coloré sémantiquement (neutre /
    orange / accent / vert) et un corps texte sur fond dégradé subtil.
    Retourne None si moins de 3 cases sur 4 sont remplies.
    """
    cases = [
        ('Contexte', contexte, C_ACCENT_MID, '01'),
        ('Problème', probleme, C_WARN,       '02'),
        ('Solution', solution, C_ACCENT,     '03'),
        ('Résultat', resultat, C_SUCCESS,    '04'),
    ]
    filled = sum(1 for _, body, _, _ in cases if body and str(body).strip())
    if filled < 3:
        return None

    slide = _content_base(prs, title, section_label)

    # ── En-tête ──────────────────────────────────────────────────────
    top_title = 0.5 if section_label else 0.18
    _icon_chip(slide, 0.22, top_title + 0.12, 0.5, CHIP_CASE,
               bg_color=C_ACCENT)
    tb_title = _tb(slide, 0.85, top_title, 12.27, 1.0)
    tf = tb_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_runs(p, title, 24, C_WHITE, base_bold=True, font=FONT_DISPLAY)
    _enable_shrink_to_fit(tf)

    sep_top = top_title + 1.0
    _rect(slide, 0.22, sep_top, 12.9, 0.025, C_ACCENT)

    # ── Grille 2x2 ───────────────────────────────────────────────────
    body_top = sep_top + 0.28
    body_bot = 7.15
    grid_left = 0.30
    grid_w    = 13.33 - 2 * grid_left
    grid_h    = body_bot - body_top
    gap       = 0.22
    card_w    = (grid_w - gap) / 2
    card_h    = (grid_h - gap) / 2

    positions = [
        (grid_left,                body_top),                  # 01
        (grid_left + card_w + gap, body_top),                  # 02
        (grid_left,                body_top + card_h + gap),   # 03
        (grid_left + card_w + gap, body_top + card_h + gap),   # 04
    ]

    band_h = 0.55  # bandeau supérieur de la carte

    for i, (label, body, color, num) in enumerate(cases):
        x, y = positions[i]

        # Carte fond + ombre portée
        card = _rounded_rect(slide, x, y, card_w, card_h, C_BG_CARD)
        _fill_gradient(card, [(0, C_ACCENT_DARK), (100, C_BG_CARD)],
                       angle_deg=135)
        _add_shadow(card, blur_pt=10, dist_pt=3, alpha_pct=40, dir_deg=90)

        # Bandeau supérieur coloré (sémantique du rôle de la carte)
        band = _rect(slide, x, y, card_w, band_h, color)
        _fill_gradient(band, [(0, color), (100, C_ACCENT_DARK)], angle_deg=0)

        # Numéro 01..04 dans le bandeau (gauche)
        tb_num = _tb(slide, x + 0.20, y, 0.95, band_h)
        tf_n = tb_num.text_frame
        tf_n.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_n.margin_left   = Inches(0)
        tf_n.margin_right  = Inches(0)
        tf_n.margin_top    = Inches(0)
        tf_n.margin_bottom = Inches(0)
        p_n = tf_n.paragraphs[0]
        run_n = p_n.add_run()
        run_n.text = num
        _run_fmt(run_n, 22, C_WHITE, bold=True, font=FONT_DISPLAY)

        # Label dans le bandeau (droite du numéro)
        tb_l = _tb(slide, x + 1.10, y, card_w - 1.20, band_h)
        tf_l = tb_l.text_frame
        tf_l.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf_l.margin_left   = Inches(0)
        tf_l.margin_right  = Inches(0.08)
        tf_l.margin_top    = Inches(0)
        tf_l.margin_bottom = Inches(0)
        tf_l.word_wrap = False
        p_l = tf_l.paragraphs[0]
        run_l = p_l.add_run()
        run_l.text = label.upper()
        _run_fmt(run_l, 14, C_WHITE, bold=True, font=FONT_DISPLAY)

        # Corps texte
        body_y    = y + band_h + 0.18
        body_h_tb = card_h - band_h - 0.30
        tb_b = _tb(slide, x + 0.22, body_y, card_w - 0.44, body_h_tb)
        tf_b = tb_b.text_frame
        tf_b.word_wrap = True
        tf_b.vertical_anchor = MSO_ANCHOR.TOP
        tf_b.margin_left   = Inches(0.04)
        tf_b.margin_right  = Inches(0.04)
        tf_b.margin_top    = Inches(0)
        tf_b.margin_bottom = Inches(0)

        if body and str(body).strip():
            text = _truncate(_clean(str(body)), 350)
            p_b = tf_b.paragraphs[0]
            _add_runs(p_b, text, 12, C_TEXT_LIGHT)
        else:
            # Carte vide : placeholder muted en italique
            p_e = tf_b.paragraphs[0]
            run_e = p_e.add_run()
            run_e.text = '— à compléter —'
            _run_fmt(run_e, 11, C_TEXT_MUTED, italic=True, font=FONT_BODY)
        # Carte de cas pratique : si le texte tronqué à 350 chars déborde
        # quand même (rare mais possible avec beaucoup de retours ligne),
        # PowerPoint réduit la police au rendu.
        _enable_shrink_to_fit(tf_b)

    return slide


def _make_timeline_slide(prs, title: str, events: list, section_label: str = ''):
    """
    Frise chronologique horizontale : 2 à 7 événements datés.

    Rendu : ligne horizontale au centre vertical, chips numérotés sur la
    ligne, dates en couleur accent2 au-dessus, titres + descriptions en
    dessous. Mise en page adaptative selon le nombre d'événements.

    events : liste de dicts {'date', 'titre', 'description'} ou de strings
    (titres seuls). Retourne None si moins de 2 événements exploitables.
    """
    # Normalisation
    norm = []
    for e in (events or []):
        if isinstance(e, dict):
            date  = _clean(str(e.get('date') or e.get('annee') or e.get('année')
                                or e.get('year') or ''))
            titre = _clean(str(e.get('titre') or e.get('title')
                                or e.get('event') or ''))
            desc  = _clean(str(e.get('description') or e.get('desc') or ''))
        else:
            date, titre, desc = '', _clean(str(e)), ''
        if date or titre or desc:
            norm.append({'date': date, 'titre': titre, 'description': desc})
    if len(norm) < 2:
        return None
    events_n = norm[:7]
    n = len(events_n)

    slide = _content_base(prs, title, section_label)

    # ── En-tête : chip flèche + titre + séparateur ─────────────────
    top_title = 0.5 if section_label else 0.18
    _icon_chip(slide, 0.22, top_title + 0.12, 0.5, CHIP_TIMELINE,
               bg_color=C_ACCENT)
    tb_title = _tb(slide, 0.85, top_title, 12.27, 1.0)
    tf = tb_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_runs(p, title, 24, C_WHITE, base_bold=True, font=FONT_DISPLAY)
    _enable_shrink_to_fit(tf)

    sep_top = top_title + 1.0
    _rect(slide, 0.22, sep_top, 12.9, 0.025, C_ACCENT)

    # ── Géométrie de la frise ──────────────────────────────────────
    body_top = sep_top + 0.30
    body_bot = 7.15

    # La ligne horizontale est positionnée à ~40% depuis le haut du body :
    # plus d'espace en bas pour titres + descriptions (qui sont plus longs).
    line_y = body_top + (body_bot - body_top) * 0.36

    # Marges latérales pour respirer en bord de slide
    side_pad   = 0.55
    line_left  = side_pad
    line_right = 13.33 - side_pad
    line_w     = line_right - line_left

    # Tailles adaptatives selon n
    if n <= 3:
        chip_d, date_pt, title_pt, desc_pt, desc_max = 0.62, 22, 17, 12, 220
    elif n == 4:
        chip_d, date_pt, title_pt, desc_pt, desc_max = 0.52, 19, 15, 11, 140
    elif n == 5:
        chip_d, date_pt, title_pt, desc_pt, desc_max = 0.45, 17, 14, 10, 100
    elif n == 6:
        chip_d, date_pt, title_pt, desc_pt, desc_max = 0.40, 15, 12, 10, 75
    else:  # 7
        chip_d, date_pt, title_pt, desc_pt, desc_max = 0.36, 13, 11, 9, 55

    # Ligne principale (rectangle fin avec dégradé violet → accent → violet)
    line_h = 0.06
    line   = _rect(slide, line_left, line_y - line_h / 2, line_w, line_h,
                   C_ACCENT_MID)
    _fill_gradient(line, [
        (0,   C_ACCENT_DARK),
        (50,  C_ACCENT),
        (100, C_ACCENT_DARK),
    ], angle_deg=0)

    # Positions horizontales des chips (équidistantes ; centré si n=1 — exclu)
    if n == 1:
        positions = [line_left + line_w / 2]
    else:
        positions = [line_left + line_w * i / (n - 1) for i in range(n)]

    # Largeur de la "colonne" allouée à chaque événement (texte au-dessus/dessous)
    col_w = line_w / n
    col_w = min(col_w, 2.6)        # plafond pour éviter du texte trop étiré
    col_w = max(col_w, 1.2)        # plancher pour rester lisible

    # ── Pour chaque événement ──────────────────────────────────────
    for i, ev in enumerate(events_n):
        cx = positions[i]

        # Date au-dessus de la ligne ─────────────────────────────────
        date_h = 0.45
        date_y = line_y - chip_d / 2 - 0.08 - date_h
        date_x = cx - col_w / 2
        tb_date = _tb(slide, date_x, date_y, col_w, date_h)
        tf_d = tb_date.text_frame
        tf_d.word_wrap = True
        tf_d.vertical_anchor = MSO_ANCHOR.BOTTOM
        tf_d.margin_left  = Inches(0.02)
        tf_d.margin_right = Inches(0.02)
        tf_d.margin_top    = Inches(0)
        tf_d.margin_bottom = Inches(0)
        p_d = tf_d.paragraphs[0]
        p_d.alignment = PP_ALIGN.CENTER
        run_d = p_d.add_run()
        run_d.text = ev['date'] or '—'
        _run_fmt(run_d, date_pt, C_ACCENT2, bold=True, font=FONT_DISPLAY)

        # Petit tiret entre la date et le chip (renforce le lien visuel)
        tick_top = date_y + date_h
        tick_h   = (line_y - chip_d / 2) - tick_top - 0.02
        if tick_h > 0:
            _rect(slide, cx - 0.012, tick_top, 0.024, tick_h, C_ACCENT_MID)

        # Chip numéroté sur la ligne ─────────────────────────────────
        _icon_chip(slide, cx - chip_d / 2, line_y - chip_d / 2, chip_d,
                   str(i + 1), bg_color=C_ACCENT,
                   glyph_size_pt=int(chip_d * 26))

        # Titre + description en dessous ─────────────────────────────
        text_top = line_y + chip_d / 2 + 0.18
        text_h   = body_bot - text_top - 0.05
        text_x   = cx - col_w / 2
        tb_t = _tb(slide, text_x, text_top, col_w, text_h)
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True
        tf_t.vertical_anchor = MSO_ANCHOR.TOP
        tf_t.margin_left  = Inches(0.04)
        tf_t.margin_right = Inches(0.04)
        tf_t.margin_top    = Inches(0)
        tf_t.margin_bottom = Inches(0)

        if ev['titre']:
            p_t = tf_t.paragraphs[0]
            p_t.alignment = PP_ALIGN.CENTER
            p_t.space_after = Pt(3)
            _add_runs(p_t, ev['titre'], title_pt, C_WHITE,
                      base_bold=True, font=FONT_DISPLAY)
        if ev['description']:
            p_dsc = tf_t.add_paragraph() if ev['titre'] else tf_t.paragraphs[0]
            p_dsc.alignment = PP_ALIGN.CENTER
            p_dsc.space_before = Pt(2)
            _add_runs(p_dsc, _truncate(ev['description'], desc_max),
                      desc_pt, C_TEXT_LIGHT)
        # Cellule timeline étroite (col_w plafonnée) : si titre+description
        # ne tiennent pas, PowerPoint réduit la police au rendu.
        _enable_shrink_to_fit(tf_t)

    return slide


def _parse_versus(text: str) -> dict | None:
    """
    Parse un bloc « versus / comparaison » en deux colonnes.

    Formes supportées :
      1. Deux sous-titres H3 (### Avant / ### Après) avec bullets sous chacun
      2. Tableau markdown à 2 colonnes (headers deviennent les labels)
      3. Paragraphe intro + deux H4/sous-listes
      4. Fallback : alterne les bullets (impair = gauche, pair = droite)
         si au moins 4 bullets sont présents

    Retourne {'left_label', 'left_items', 'right_label', 'right_items'}
    ou None si la structure n'est pas exploitable (< 2 items de chaque côté).
    """
    # ── Cas 1 : deux H3 ──────────────────────────────────────────────
    h3_re = re.compile(r'^#{3,4}\s+(.+)$', re.MULTILINE)
    h3_matches = list(h3_re.finditer(text))
    if len(h3_matches) >= 2:
        blocks = []
        for i, m in enumerate(h3_matches[:2]):
            label = _clean(m.group(1))
            start = m.end()
            end   = h3_matches[i + 1].start() if i + 1 < len(h3_matches) else len(text)
            body  = text[start:end]
            items = _extract_bullets(body, max_b=6)
            if not items:
                # Fallback : phrases courtes du paragraphe
                para = _first_paragraph(body)
                if para:
                    items = [s.strip() for s in para.split('.') if s.strip()][:6]
            blocks.append((label, items))
        if len(blocks[0][1]) >= 1 and len(blocks[1][1]) >= 1:
            return {
                'left_label': blocks[0][0], 'left_items': blocks[0][1],
                'right_label': blocks[1][0], 'right_items': blocks[1][1],
            }

    # ── Cas 2 : tableau markdown 2 colonnes ──────────────────────────
    headers, rows = _parse_md_table(text)
    if len(headers) == 2 and len(rows) >= 2:
        left_items  = [r[0] for r in rows if len(r) > 0 and r[0].strip()][:6]
        right_items = [r[1] for r in rows if len(r) > 1 and r[1].strip()][:6]
        if len(left_items) >= 1 and len(right_items) >= 1:
            return {
                'left_label': _clean(headers[0]),
                'left_items': left_items,
                'right_label': _clean(headers[1]),
                'right_items': right_items,
            }

    # ── Cas 3 : parser « Avant : ... / Après : ... » inline ──────────
    # Forme : "**Avant** : liste1 ; item2 ; item3" et "**Après** : ..."
    inline = re.findall(r'\*\*([^*]+?)\*\*\s*[:：]\s*([^\n]+)', text)
    if len(inline) >= 2:
        left_items  = [s.strip() for s in re.split(r'[;•]', inline[0][1]) if s.strip()][:6]
        right_items = [s.strip() for s in re.split(r'[;•]', inline[1][1]) if s.strip()][:6]
        if len(left_items) >= 2 and len(right_items) >= 2:
            return {
                'left_label': _clean(inline[0][0]),
                'left_items': left_items,
                'right_label': _clean(inline[1][0]),
                'right_items': right_items,
            }

    return None


def _is_pro_con_label(label: str) -> str:
    """
    Classifie un label de colonne versus :
    'pro'  → avantages / après / moderne / ...  (vert)
    'con'  → inconvénients / avant / traditionnel / ...  (orange)
    'neutral' → n'appartient ni à l'un ni à l'autre
    """
    lbl = label.lower().strip()
    pro_kw = ('après', 'apres', 'pour', 'avantages', 'bénéfices', 'benefices',
              'forces', 'moderne', 'nouveau', 'positif', 'pros', 'pratique')
    con_kw = ('avant', 'contre', 'inconvénients', 'inconvenients', 'faiblesses',
              'traditionnel', 'ancien', 'négatif', 'negatif', 'cons', 'théorie',
              'theorie', 'risques', 'limites')
    if any(k in lbl for k in pro_kw):
        return 'pro'
    if any(k in lbl for k in con_kw):
        return 'con'
    return 'neutral'


def _make_versus_slide(prs, title: str, left_label: str, left_items: list[str],
                       right_label: str, right_items: list[str],
                       section_label: str = ''):
    """
    Slide « versus / comparaison » deux colonnes.

    Chaque colonne a son en-tête coloré distinctif (orange vs vert par défaut ;
    accent violet si les labels sont neutres) et affiche 2 à 6 items.
    Les items côté « con » sont marqués d'une croix ✗, ceux côté « pro » d'un ✓.
    Un séparateur central avec badge « VS » marque la dualité.
    """
    left_items  = [_clean(str(x)) for x in (left_items or []) if str(x).strip()][:6]
    right_items = [_clean(str(x)) for x in (right_items or []) if str(x).strip()][:6]
    if not left_items or not right_items:
        return None

    slide = _content_base(prs, title, section_label)

    # ── En-tête ──────────────────────────────────────────────────────
    top_title = 0.5 if section_label else 0.18
    _icon_chip(slide, 0.22, top_title + 0.12, 0.5, CHIP_COMPARE, bg_color=C_ACCENT)
    tb_title = _tb(slide, 0.85, top_title, 12.27, 1.0)
    tf = tb_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_runs(p, title, 24, C_WHITE, base_bold=True, font=FONT_DISPLAY)
    _enable_shrink_to_fit(tf)

    sep_top = top_title + 1.0
    _rect(slide, 0.22, sep_top, 12.9, 0.025, C_ACCENT)

    # ── Couleurs de colonnes selon classification du label ───────────
    left_kind  = _is_pro_con_label(left_label)
    right_kind = _is_pro_con_label(right_label)
    # Si on a exactement un pro et un con, on applique les couleurs sémantiques
    if {left_kind, right_kind} == {'pro', 'con'}:
        left_color  = C_WARN    if left_kind == 'con' else C_SUCCESS
        right_color = C_SUCCESS if right_kind == 'pro' else C_WARN
        left_bullet  = '✗' if left_kind == 'con' else '✓'
        right_bullet = '✓' if right_kind == 'pro' else '✗'
    else:
        # Labels neutres : on reste sur la palette accent (pas de sémantique)
        left_color  = C_ACCENT
        right_color = C_ACCENT2
        left_bullet  = '◆'
        right_bullet = '◆'

    # ── Zone deux colonnes ───────────────────────────────────────────
    body_top = sep_top + 0.35
    body_bot = 7.15
    body_h   = body_bot - body_top

    col_gap     = 0.50
    total_w     = 13.33 - 0.44 - col_gap          # après side-pad 0.22 * 2
    col_w       = total_w / 2
    left_x      = 0.22
    right_x     = left_x + col_w + col_gap
    center_x    = left_x + col_w + col_gap / 2

    header_h    = 0.60

    def _render_column(x, label, items, color, bullet):
        # Bandeau titre coloré
        hdr = _rect(slide, x, body_top, col_w, header_h, color)
        _fill_gradient(hdr, [(0, color), (100, C_ACCENT_DARK)], angle_deg=0)
        _add_shadow(hdr, blur_pt=8, dist_pt=2, alpha_pct=40, dir_deg=90)
        tb = _tb(slide, x, body_top, col_w, header_h)
        tfh = tb.text_frame
        tfh.vertical_anchor = MSO_ANCHOR.MIDDLE
        tfh.margin_left = Inches(0.15)
        tfh.margin_right = Inches(0.15)
        ph = tfh.paragraphs[0]
        ph.alignment = PP_ALIGN.CENTER
        run = ph.add_run()
        run.text = (label or '').upper()
        _run_fmt(run, 16, C_WHITE, bold=True, font=FONT_DISPLAY)

        # Liste des items
        items_top = body_top + header_h + 0.20
        items_h   = body_h - header_h - 0.20
        n = len(items)
        # Police adaptative : 2-3 → 16pt ; 4 → 14pt ; 5-6 → 12pt
        if n <= 3:
            item_pt, lead_pt = 16, 17
            item_gap = 0.18
        elif n == 4:
            item_pt, lead_pt = 14, 15
            item_gap = 0.14
        else:
            item_pt, lead_pt = 12, 13
            item_gap = 0.10

        row_h = (items_h - item_gap * (n - 1)) / n
        for i, it in enumerate(items):
            y = items_top + i * (row_h + item_gap)
            # Carte fond subtile par item (améliore la lisibilité)
            card = _rounded_rect(slide, x + 0.05, y, col_w - 0.10, row_h, C_BG_CARD)
            _fill_gradient(card, [(0, C_ACCENT_DARK), (100, C_BG_CARD)], angle_deg=0)

            # Chip bullet (coche ou croix selon pro/con)
            chip_d = min(0.40, row_h - 0.10)
            _icon_chip(slide, x + 0.15, y + (row_h - chip_d) / 2, chip_d,
                       bullet, bg_color=color, glyph_size_pt=int(chip_d * 22))

            # Texte item — lead:corps détecté si présent
            text_left = x + 0.20 + chip_d + 0.20
            text_w    = col_w - (text_left - x) - 0.15
            tb = _tb(slide, text_left, y, text_w, row_h)
            tfi = tb.text_frame
            tfi.word_wrap = True
            tfi.vertical_anchor = MSO_ANCHOR.MIDDLE
            tfi.margin_left = Inches(0.04)
            tfi.margin_right = Inches(0.04)
            tfi.margin_top = Inches(0)
            tfi.margin_bottom = Inches(0)
            p_i = tfi.paragraphs[0]

            lead, rest = '', it
            mb = re.match(r'^\*\*(.+?)\*\*\s*[:：]?\s*(.*)', it)
            if mb and mb.group(2).strip():
                lead = _clean(mb.group(1))
                rest = _clean(mb.group(2))

            if lead:
                rv = p_i.add_run()
                rv.text = lead + ' '
                _run_fmt(rv, lead_pt, C_WHITE, bold=True, font=FONT_DISPLAY)
                _add_runs(p_i, rest, item_pt, C_TEXT_LIGHT)
            else:
                _add_runs(p_i, it, item_pt, C_TEXT_LIGHT, base_bold=False)
            _enable_shrink_to_fit(tfi)

    _render_column(left_x, left_label, left_items, left_color, left_bullet)
    _render_column(right_x, right_label, right_items, right_color, right_bullet)

    # ── Séparateur central + badge « VS » ────────────────────────────
    vs_d = 0.70
    vs_y = body_top + (body_h - vs_d) / 2
    vs_x = center_x - vs_d / 2

    # Trait pointillé haut + bas autour du badge
    _rect(slide, center_x - 0.01, body_top + header_h + 0.15,
          0.02, vs_y - body_top - header_h - 0.15, C_ACCENT_MID)
    _rect(slide, center_x - 0.01, vs_y + vs_d,
          0.02, body_bot - vs_y - vs_d - 0.05, C_ACCENT_MID)

    # Badge circulaire « VS »
    vs_chip = _rounded_rect(slide, vs_x, vs_y, vs_d, vs_d, C_ACCENT)
    _fill_gradient(vs_chip, [(0, C_ACCENT), (100, C_ACCENT_DARK)], angle_deg=135)
    _add_shadow(vs_chip, blur_pt=10, dist_pt=3, alpha_pct=50, dir_deg=90)
    tb_vs = _tb(slide, vs_x, vs_y, vs_d, vs_d)
    tfv = tb_vs.text_frame
    tfv.vertical_anchor = MSO_ANCHOR.MIDDLE
    tfv.margin_left = Inches(0)
    tfv.margin_right = Inches(0)
    pv = tfv.paragraphs[0]
    pv.alignment = PP_ALIGN.CENTER
    rv = pv.add_run()
    rv.text = 'VS'
    _run_fmt(rv, 20, C_WHITE, bold=True, font=FONT_DISPLAY)

    return slide


# ═══════════════════════════════════════════════════════
#  SLIDES V2 — layouts du pipeline multi-agents
# ═══════════════════════════════════════════════════════

def _parse_numeric(value_str) -> tuple:
    """
    Extrait (nombre, suffixe) d'une chaîne. Exemples :
    "73%"        -> (73.0, "%")
    "1 500 €"    -> (1500.0, "€")
    "élevé"      -> (None, "")
    "2,5 M€"     -> (2.5, "M€")
    """
    s = str(value_str).strip()
    m = re.search(r'-?\d+(?:[\s\u202f]\d+)*(?:[.,]\d+)?', s)
    if not m:
        return None, ''
    raw = m.group(0).replace(' ', '').replace('\u202f', '').replace(',', '.')
    try:
        num = float(raw)
    except ValueError:
        return None, ''
    suffix = s[m.end():].strip()[:4]
    return num, suffix


def _chart_area_transparent(chart):
    """
    Force le fond du chart-space et du plot area à transparent,
    pour que le chart s'intègre visuellement au fond sombre de la slide.
    """
    chart_space = chart._chartSpace
    # 1) chartSpace > spPr
    spPr = chart_space.find(qn('c:spPr'))
    if spPr is None:
        spPr = etree.SubElement(chart_space, qn('c:spPr'))
    for child in list(spPr):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('solidFill', 'gradFill', 'pattFill', 'blipFill', 'noFill', 'ln'):
            spPr.remove(child)
    etree.SubElement(spPr, f'{{{_DML_NS}}}noFill')
    ln = etree.SubElement(spPr, f'{{{_DML_NS}}}ln')
    etree.SubElement(ln, f'{{{_DML_NS}}}noFill')

    # 2) plotArea > spPr
    plotArea = chart_space.find('.//' + qn('c:plotArea'))
    if plotArea is not None:
        plotSpPr = plotArea.find(qn('c:spPr'))
        if plotSpPr is None:
            plotSpPr = etree.SubElement(plotArea, qn('c:spPr'))
        for child in list(plotSpPr):
            tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
            if tag in ('solidFill', 'gradFill', 'pattFill', 'blipFill', 'noFill', 'ln'):
                plotSpPr.remove(child)
        etree.SubElement(plotSpPr, f'{{{_DML_NS}}}noFill')


def _make_progress_slide(prs, title: str, stats: list, section_label: str = ''):
    """
    Slide stats en barres de progression horizontales (track + remplissage).
    Déclenchée quand 1 à 4 stats sont TOUTES en pourcentage (0-100%).
    Retourne None sinon — le caller doit alors essayer un autre rendu
    (chart natif ou cartes texte).

    Lecture visuelle immédiate : chaque barre montre "où on en est sur 100".
    Plus expressif qu'un bar chart pour les %, car la barre évoque la jauge
    « remplie / vide ».
    """
    if not stats:
        return None

    parsed = []
    for stat in stats[:4]:
        num, suffix = _parse_numeric(stat.get('valeur', ''))
        if num is None or suffix != '%' or not (0 <= num <= 100):
            # Règle stricte : on ne mélange pas des stats % et non-%,
            # sinon la métaphore « barre 0-100 » n'a plus de sens.
            return None
        parsed.append({'num': num, 'label': str(stat.get('label', ''))})

    if not parsed:
        return None

    n = len(parsed)
    slide = _content_base(prs, title, section_label)

    # En-tête (titre + trait accent) — aligné sur les autres slides stat/chart
    top_title = 0.5 if section_label else 0.18
    _icon_chip(slide, 0.22, top_title + 0.12, 0.5, CHIP_STATS, bg_color=C_ACCENT)
    tb_title = _tb(slide, 0.85, top_title, 12.27, 1.0)
    tf = tb_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_runs(p, title, 24, C_WHITE, base_bold=True, font=FONT_DISPLAY)
    _enable_shrink_to_fit(tf)

    sep_top = top_title + 1.0
    _rect(slide, 0.22, sep_top, 12.9, 0.025, C_ACCENT)

    # Zone de contenu : entre le trait accent et la zone pied de page.
    # Footer est à y=7.30 → on se limite à 7.15 pour respirer.
    body_top = sep_top + 0.45
    body_bot = 7.15
    body_h   = body_bot - body_top
    row_h    = body_h / n

    # Tailles adaptatives (plus n est petit, plus on grossit)
    if n == 1:
        label_pt, value_pt, bar_h = 18, 52, 0.50
    elif n == 2:
        label_pt, value_pt, bar_h = 16, 40, 0.42
    else:
        label_pt, value_pt, bar_h = 14, 28, 0.32

    # Géométrie horizontale
    side_pad   = 0.22
    track_left = side_pad
    track_w    = 13.33 - 2 * side_pad   # ~12.89"

    for i, item in enumerate(parsed):
        num   = item['num']
        label = item['label']

        row_center = body_top + (i + 0.5) * row_h
        # On place : label + valeur sur une ligne, puis barre juste dessous.
        # Le couple (texte / barre) est centré verticalement dans la row.
        text_h = max(label_pt, value_pt) / 72.0 + 0.10  # hauteur approx. en pouces
        block_h = text_h + 0.12 + bar_h
        block_top = row_center - block_h / 2

        # --- Label à gauche ---
        tb_lbl = _tb(slide, track_left, block_top, track_w * 0.72, text_h)
        tf_l = tb_lbl.text_frame
        tf_l.word_wrap = True
        pl = tf_l.paragraphs[0]
        pl.alignment = PP_ALIGN.LEFT
        _add_runs(pl, _truncate(label, 80), label_pt, C_TEXT_LIGHT,
                  base_bold=True, font=FONT_BODY)

        # --- Valeur à droite (aligné à la fin de la barre) ---
        val_w = track_w * 0.26
        val_left = track_left + track_w - val_w
        tb_val = _tb(slide, val_left, block_top, val_w, text_h)
        tf_v = tb_val.text_frame
        pv = tf_v.paragraphs[0]
        pv.alignment = PP_ALIGN.RIGHT
        val_text = f"{num:g}%"
        run_v = pv.add_run()
        run_v.text = val_text
        _run_fmt(run_v, value_pt, C_ACCENT2, bold=True, font=FONT_DISPLAY)

        # --- Track (fond de la barre, pleine largeur) ---
        bar_top = block_top + text_h + 0.12
        track = _rounded_rect(slide, track_left, bar_top, track_w, bar_h,
                              C_ACCENT_DARK)
        _fill_gradient(track, [(0, C_ACCENT_DARK), (100, C_ACCENT_MID)],
                       angle_deg=90)

        # --- Fill (remplissage proportionnel) ---
        # Minimum visible pour qu'à 0% on voie tout de même un embryon de barre.
        fill_w = max(track_w * (num / 100.0), 0.05) if num > 0 else 0
        if fill_w > 0:
            fill = _rounded_rect(slide, track_left, bar_top, fill_w, bar_h,
                                 C_ACCENT)
            _fill_gradient(fill, [(0, C_ACCENT), (100, C_ACCENT2)],
                           angle_deg=0)
            _add_shadow(fill, blur_pt=8, dist_pt=2, alpha_pct=45, dir_deg=90)

    return slide


def _make_stat_chart_slide(prs, title: str, stats: list, section_label: str = ''):
    """
    Slide stats avec graphique natif (bar chart horizontal) sur fond sombre.
    Retourne None si moins de 2 valeurs numériques — le caller doit alors
    retomber sur _make_stat_slide (cartes texte).
    """
    parsed = []
    for stat in stats[:6]:
        num, suffix = _parse_numeric(stat.get('valeur', ''))
        if num is None:
            continue
        parsed.append({
            'num': num,
            'suffix': suffix,
            'label': str(stat.get('label', '')),
        })
    if len(parsed) < 2:
        return None

    slide = _content_base(prs, title, section_label)

    top_title = 0.5 if section_label else 0.18
    _icon_chip(slide, 0.22, top_title + 0.12, 0.5, CHIP_STATS, bg_color=C_ACCENT)
    tb_title = _tb(slide, 0.85, top_title, 12.27, 1.0)
    tf = tb_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_runs(p, title, 24, C_WHITE, base_bold=True, font=FONT_DISPLAY)
    _enable_shrink_to_fit(tf)

    sep_top = top_title + 1.0
    _rect(slide, 0.22, sep_top, 12.9, 0.025, C_ACCENT)

    # BAR_CLUSTERED affiche la 1ère catégorie en bas : on inverse pour afficher
    # le 1er stat en haut (lecture naturelle).
    ordered = list(reversed(parsed))
    chart_data = CategoryChartData()
    chart_data.categories = [_truncate(p['label'], 45) for p in ordered]
    chart_data.add_series('Valeur', [p['num'] for p in ordered])

    # Format numérique : % si tous les suffixes sont %, sinon nombre brut
    suffixes = {p['suffix'] for p in parsed}
    if suffixes == {'%'}:
        num_format = '0"%"'
    else:
        num_format = '0'

    chart_top = sep_top + 0.35
    chart_h   = 7.5 - chart_top - 0.55   # marge pour le footer
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.35), Inches(chart_top),
        Inches(12.6), Inches(chart_h),
        chart_data,
    )
    chart = gframe.chart

    chart.has_title  = False
    chart.has_legend = False

    # Labels de données en bout de barre
    plot = chart.plots[0]
    plot.has_data_labels = True
    dlbl = plot.data_labels
    dlbl.show_value = True
    dlbl.number_format = num_format
    dlbl.position = XL_LABEL_POSITION.OUTSIDE_END
    dlbl.font.size = Pt(14)
    dlbl.font.bold = True
    dlbl.font.color.rgb = C_WHITE
    dlbl.font.name = FONT_DISPLAY

    # Couleur des barres (accent violet) + pas de bordure
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = C_ACCENT
    series.format.line.fill.background()

    # Axe catégories (labels textuels à gauche des barres)
    try:
        cat_ax = chart.category_axis
        cat_ax.tick_labels.font.size = Pt(12)
        cat_ax.tick_labels.font.color.rgb = C_TEXT_LIGHT
        cat_ax.tick_labels.font.name = FONT_BODY
        cat_ax.format.line.fill.background()
    except Exception:
        pass

    # Axe des valeurs (bas) : masqué
    try:
        chart.value_axis.visible = False
    except Exception:
        pass

    # Fond chart+plot transparent pour se fondre au slide
    _chart_area_transparent(chart)

    return slide


def _make_stat_slide(prs, title: str, stats: list[dict], section_label: str = ''):
    """Slide stat-callout : statistiques en grands nombres sur fond de cartes."""
    slide = _content_base(prs, title, section_label)

    top_title = 0.5 if section_label else 0.18
    _icon_chip(slide, 0.22, top_title + 0.12, 0.5, CHIP_STATS, bg_color=C_ACCENT)
    tb_title = _tb(slide, 0.85, top_title, 12.27, 1.0)
    tf = tb_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_runs(p, title, 24, C_WHITE, base_bold=True, font=FONT_DISPLAY)
    _enable_shrink_to_fit(tf)

    sep_top = top_title + 1.0
    _rect(slide, 0.22, sep_top, 12.9, 0.025, C_ACCENT)
    body_top = sep_top + 0.5

    n = min(len(stats), 4)
    if n == 0:
        return slide
    card_w = 12.9 / n

    for i, stat in enumerate(stats[:4]):
        lx = 0.22 + i * card_w
        card = _rounded_rect(slide, lx + 0.12, body_top, card_w - 0.25, 3.8, C_ACCENT_DARK)
        _fill_gradient(card, [(0, C_ACCENT_MID), (100, C_ACCENT_DARK)], angle_deg=135)
        _add_shadow(card, blur_pt=12, dist_pt=4, alpha_pct=50, dir_deg=90)

        # Valeur
        tb_val = _tb(slide, lx + 0.12, body_top + 0.5, card_w - 0.25, 1.8)
        tf_v = tb_val.text_frame
        pv = tf_v.paragraphs[0]
        pv.alignment = PP_ALIGN.CENTER
        run_v = pv.add_run()
        run_v.text = str(stat.get('valeur', ''))
        _run_fmt(run_v, 44, C_ACCENT2, bold=True, font=FONT_DISPLAY)

        # Label
        tb_lbl = _tb(slide, lx + 0.12, body_top + 2.5, card_w - 0.25, 1.1)
        tf_l = tb_lbl.text_frame
        tf_l.word_wrap = True
        pl = tf_l.paragraphs[0]
        pl.alignment = PP_ALIGN.CENTER
        _add_runs(pl, _truncate(str(stat.get('label', '')), 80), 12, C_TEXT_MUTED)

    return slide


def _make_two_column_text_slide(prs, title: str, left_text: str, right_text: str,
                                 section_label: str = ''):
    """Slide two-column avec texte libre gauche/droite (pas une liste de puces)."""
    slide = _content_base(prs, title, section_label)

    top_title = 0.5 if section_label else 0.18
    _icon_chip(slide, 0.22, top_title + 0.12, 0.5, CHIP_COMPARE, bg_color=C_ACCENT)
    tb_title = _tb(slide, 0.85, top_title, 12.27, 1.0)
    tf = tb_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_runs(p, title, 24, C_WHITE, base_bold=True, font=FONT_DISPLAY)
    _enable_shrink_to_fit(tf)

    sep_top = top_title + 1.0
    _rect(slide, 0.22, sep_top, 12.9, 0.025, C_ACCENT)
    body_top = sep_top + 0.18
    body_h   = 7.5 - body_top - 0.35

    _rect(slide, 6.8, body_top, 0.02, body_h, C_ACCENT_MID)

    for text, lx in [(left_text, 0.22), (right_text, 6.95)]:
        tb = _tb(slide, lx, body_top, 6.35, body_h)
        tf = tb.text_frame
        tf.word_wrap = True
        lines = [l for l in text.split('\n') if l.strip()]
        for i, line in enumerate(lines[:14]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(4)
            _add_runs(p, line.strip(), 12, C_TEXT_LIGHT)
        _enable_shrink_to_fit(tf)

    return slide


def _make_stepper_slide(prs, title: str, steps: list, section_label: str = ''):
    """
    Stepper vertical pour processus séquentiels (2 à 6 étapes).

    Chaque étape est représentée par un cercle numéroté (chip) relié au suivant
    par une ligne verticale, avec un titre et une description à droite.
    Mise en page adaptative selon le nombre d'étapes.

    steps : liste de dicts {'titre', 'description'} ou de strings (titres seuls).
    Retourne None si moins de 2 étapes exploitables.
    """
    # Normalisation
    norm = []
    for s in (steps or []):
        if isinstance(s, dict):
            titre = _clean(str(s.get('titre') or s.get('title') or s.get('nom') or ''))
            desc  = _clean(str(s.get('description') or s.get('desc') or s.get('contenu') or ''))
        else:
            titre, desc = _clean(str(s)), ''
        if titre or desc:
            norm.append({'titre': titre, 'description': desc})
    if len(norm) < 2:
        return None
    steps_n = norm[:6]
    n = len(steps_n)

    slide = _content_base(prs, title, section_label)

    # ── En-tête : chip stepper + titre + séparateur
    top_title = 0.5 if section_label else 0.18
    _icon_chip(slide, 0.22, top_title + 0.12, 0.5, CHIP_STEPS, bg_color=C_ACCENT)
    tb_title = _tb(slide, 0.85, top_title, 12.27, 1.0)
    tf = tb_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_runs(p, title, 24, C_WHITE, base_bold=True, font=FONT_DISPLAY)
    _enable_shrink_to_fit(tf)

    sep_top = top_title + 1.0
    _rect(slide, 0.22, sep_top, 12.9, 0.025, C_ACCENT)

    # ── Zone stepper
    body_top = sep_top + 0.30
    body_bot = 7.15
    body_h   = body_bot - body_top

    # Dimensions adaptatives : plus d'étapes = plus petites
    if n <= 3:
        chip_d, num_pt, title_pt, desc_pt, desc_max = 0.85, 26, 18, 13, 220
    elif n == 4:
        chip_d, num_pt, title_pt, desc_pt, desc_max = 0.75, 22, 16, 12, 160
    elif n == 5:
        chip_d, num_pt, title_pt, desc_pt, desc_max = 0.65, 19, 15, 11, 120
    else:  # 6
        chip_d, num_pt, title_pt, desc_pt, desc_max = 0.55, 17, 14, 10, 90

    step_h        = body_h / n
    chip_left     = 0.60
    chip_center_x = chip_left + chip_d / 2
    connector_w   = 0.05
    text_left     = chip_left + chip_d + 0.35
    text_w        = 13.33 - text_left - 0.40

    for i, step in enumerate(steps_n):
        y_center = body_top + step_h * (i + 0.5)
        y_chip   = y_center - chip_d / 2

        # Connecteur vertical depuis l'étape précédente
        if i > 0:
            prev_y_center = body_top + step_h * (i - 0.5)
            conn_top = prev_y_center + chip_d / 2
            conn_h   = (y_center - chip_d / 2) - conn_top
            if conn_h > 0:
                _rect(slide, chip_center_x - connector_w / 2, conn_top,
                      connector_w, conn_h, C_ACCENT_MID)

        # Chip numéroté
        _icon_chip(slide, chip_left, y_chip, chip_d, str(i + 1),
                   bg_color=C_ACCENT, glyph_size_pt=num_pt)

        # Texte : titre (+ description)
        text_h = step_h - 0.10
        tb_step = _tb(slide, text_left, y_chip - 0.05, text_w, text_h)
        tf_step = tb_step.text_frame
        tf_step.word_wrap = True
        tf_step.margin_left   = Inches(0.04)
        tf_step.margin_right  = Inches(0.04)
        tf_step.margin_top    = Inches(0)
        tf_step.margin_bottom = Inches(0)
        tf_step.vertical_anchor = MSO_ANCHOR.MIDDLE

        titre_step = step['titre']
        desc_step  = _truncate(step['description'], desc_max) if step['description'] else ''

        if titre_step:
            p_t = tf_step.paragraphs[0]
            p_t.space_after = Pt(3)
            _add_runs(p_t, titre_step, title_pt, C_WHITE, base_bold=True, font=FONT_DISPLAY)
            if desc_step:
                p_d = tf_step.add_paragraph()
                p_d.space_before = Pt(2)
                _add_runs(p_d, desc_step, desc_pt, C_TEXT_LIGHT)
        elif desc_step:
            p_d = tf_step.paragraphs[0]
            _add_runs(p_d, desc_step, desc_pt, C_TEXT_LIGHT)
        _enable_shrink_to_fit(tf_step)

    return slide


def _make_schema_slide(prs, title: str, description: str, elements: list[str],
                        section_label: str = ''):
    """Slide schema : description + éléments reliés par des flèches."""
    slide = _content_base(prs, title, section_label)

    top_title = 0.5 if section_label else 0.18
    _icon_chip(slide, 0.22, top_title + 0.12, 0.5, CHIP_SCHEMA, bg_color=C_ACCENT)
    tb_title = _tb(slide, 0.85, top_title, 12.27, 1.0)
    tf = tb_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    _add_runs(p, title, 24, C_WHITE, base_bold=True, font=FONT_DISPLAY)
    _enable_shrink_to_fit(tf)

    sep_top = top_title + 1.0
    _rect(slide, 0.22, sep_top, 12.9, 0.025, C_ACCENT)
    body_top = sep_top + 0.3

    if description:
        tb_desc = _tb(slide, 0.22, body_top, 12.9, 0.75)
        tf_d = tb_desc.text_frame
        tf_d.word_wrap = True
        pd = tf_d.paragraphs[0]
        _add_runs(pd, _truncate(description, 200), 12, C_TEXT_MUTED)
        _enable_shrink_to_fit(tf_d)
        body_top += 0.9

    n = min(len(elements), 7)
    if n == 0:
        return slide

    elem_w   = min(2.0, 12.0 / n)
    arrow_w  = 0.35
    total_w  = n * elem_w + (n - 1) * arrow_w
    start_x  = 0.22 + max(0.0, (12.9 - total_w) / 2)
    elem_top = body_top + 0.5

    for i, el in enumerate(elements[:7]):
        lx = start_x + i * (elem_w + arrow_w)
        box = _rounded_rect(slide, lx, elem_top, elem_w, 1.0, C_ACCENT_MID)
        _fill_gradient(box, [(0, C_ACCENT_MID), (100, C_ACCENT_DARK)], angle_deg=90)
        _add_shadow(box, blur_pt=10, dist_pt=3, alpha_pct=45, dir_deg=90)
        tb = _tb(slide, lx, elem_top, elem_w, 1.0)
        tf = tb.text_frame
        tf.word_wrap = True
        pe = tf.paragraphs[0]
        pe.alignment = PP_ALIGN.CENTER
        _add_runs(pe, _truncate(el, 35), 11, C_TEXT_LIGHT)
        _enable_shrink_to_fit(tf)

        if i < n - 1:
            tb_arr = _tb(slide, lx + elem_w + 0.02, elem_top + 0.28, arrow_w - 0.04, 0.45)
            tf_a = tb_arr.text_frame
            pa = tf_a.paragraphs[0]
            pa.alignment = PP_ALIGN.CENTER
            run_a = pa.add_run()
            run_a.text = '→'
            _run_fmt(run_a, 18, C_ACCENT2, bold=True, font=FONT_DISPLAY)

    return slide


# ── Point d'entrée V2 ──────────────────────────────────────────────────────

def slides_json_to_pptx(slides_json: dict, specialite: str, module: str,
                         chapitre: str, niveau: str = '',
                         title_image: bytes = None, photographer: str = '') -> bytes:
    """
    Génère une présentation PPTX depuis le slides_json de l'Agent Designer.
    Chaque slide est rendue selon son layout : bullets, two-column, stat-callout, schema.
    Utilise les mêmes styles visuels que markdown_to_pptx.
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = slides_json.get('slides', [])
    title_done = False

    # Sommaire cliquable : placeholder créé juste après la slide titre.
    # Rempli à la fin une fois les slides de section référencées disponibles.
    toc_slide = None
    toc_entries = []  # [(titre_section, slide_section), ...]

    section_counter_v2 = 0
    section_slide_indices = set()  # indices à exclure du footer (full-bleed)

    for slide_data in slides:
        slide_type = slide_data.get('type', 'content')
        layout     = slide_data.get('layout', 'bullets')
        titre      = _clean(slide_data.get('titre', ''))
        contenu    = slide_data.get('contenu') or {}

        # ── Slide titre ──────────────────────────────────────────────────────
        if slide_type == 'title' or not title_done:
            _make_title_slide(prs, specialite, module, chapitre, niveau,
                              title_image=title_image, photographer=photographer)
            title_done = True
            # Crée le placeholder du sommaire juste après le titre
            if toc_slide is None:
                toc_slide = _make_toc_placeholder(prs)
            if slide_type == 'title':
                continue

        # ── Slide section ────────────────────────────────────────────────────
        if slide_type == 'section':
            section_counter_v2 += 1
            _make_section_slide(prs, titre, numero=section_counter_v2)
            section_slide_indices.add(len(prs.slides) - 1)
            toc_entries.append((titre, prs.slides[-1]))
            continue

        # ── Layouts spécialisés détectés par mot-clé du titre (parité V1) ───
        # Route vers _make_definitions_slide / _make_key_points_slide avec
        # fallback sur le dispatch normal si le contenu n'est pas exploitable.
        titre_lower = titre.lower()
        SEC_DEF_KW  = ('définitions des concepts clés', 'définitions', 'concepts clés')
        SEC_PTS_KW  = ('points importants à retenir', 'points importants',
                       'points clés')
        SEC_GOAL_KW = ('objectifs pédagogiques', 'objectifs du cours',
                       'objectifs du chapitre', 'objectifs d\'apprentissage',
                       'compétences visées')
        SEC_RECAP_KW = ('synthèse', 'à retenir', 'récapitulatif',
                        'take-aways', 'en résumé', 'l\'essentiel')

        if layout in ('bullets', 'two-column'):
            items_raw = contenu.get('items') or []
            if not items_raw:
                left  = str(contenu.get('colonne_gauche', '') or '')
                right = str(contenu.get('colonne_droite', '') or '')
                items_raw = _extract_bullets(left + '\n' + right, 20)

            if items_raw and any(kw in titre_lower for kw in SEC_DEF_KW):
                full_text = '\n'.join(f'- {it}' for it in items_raw)
                defs = _parse_definitions(full_text)
                if defs:
                    _make_definitions_slide(prs, defs)
                    continue
                # Sinon : on laisse tomber et on fallback sur le dispatch normal

            if items_raw and any(kw in titre_lower for kw in SEC_PTS_KW):
                points = [_clean(str(it)) for it in items_raw if str(it).strip()]
                if points:
                    _make_key_points_slide(prs, points[:12])
                    continue

            if items_raw and any(kw in titre_lower for kw in SEC_GOAL_KW):
                objs = [_clean(str(it)) for it in items_raw if str(it).strip()]
                if objs and _make_objectives_slide(prs, objs) is not None:
                    continue

            if items_raw and any(kw in titre_lower for kw in SEC_RECAP_KW):
                pts = [_clean(str(it)) for it in items_raw if str(it).strip()]
                if pts and _make_synthese_slide(prs, pts) is not None:
                    continue

        # ── Slides de contenu selon layout ──────────────────────────────────
        if layout == 'bullets':
            items = contenu.get('items', [])
            if not items:
                continue
            if len(items) > MAX_BULLETS_SINGLE:
                _make_two_column_slide(prs, titre, [str(it) for it in items])
            else:
                _make_content_slide(prs, titre, [str(it) for it in items], is_bullets=True)

        elif layout == 'two-column':
            left  = str(contenu.get('colonne_gauche', '') or '')
            right = str(contenu.get('colonne_droite', '') or '')
            # Si les colonnes contiennent des listes → utilise le layout bullets existant
            left_b  = _extract_bullets(left, 20)
            right_b = _extract_bullets(right, 20)
            if left_b and right_b:
                _make_two_column_slide(prs, titre, left_b + right_b)
            else:
                _make_two_column_text_slide(prs, titre, left or '—', right or '—')

        elif layout == 'stat-callout':
            stats = contenu.get('stats', [])
            if stats:
                # Priorité : 1) barres de progression si toutes les valeurs sont
                # des % entre 0 et 100 (métaphore jauge), 2) bar chart natif si
                # ≥ 2 numériques, 3) cartes texte en dernier recours.
                if _make_progress_slide(prs, titre, stats) is None:
                    if _make_stat_chart_slide(prs, titre, stats) is None:
                        _make_stat_slide(prs, titre, stats)

        elif layout == 'schema':
            description = str(contenu.get('description_schema', '') or '')
            elements    = [str(e) for e in (contenu.get('elements') or [])]
            _make_schema_slide(prs, titre, description, elements)

        elif layout == 'callout':
            # Callout / citation : un texte court mis en exergue.
            # Champs tolérés : 'texte' (principal), 'citation', 'quote' ;
            # attribution dans 'attribution', 'auteur' ou 'source'.
            quote = (contenu.get('texte') or contenu.get('citation')
                     or contenu.get('quote') or '')
            attribution = (contenu.get('attribution') or contenu.get('auteur')
                           or contenu.get('source') or '')
            quote = str(quote).strip()
            attribution = str(attribution).strip()
            if quote:
                _make_callout_slide(prs, quote, attribution, section_label=titre)

        elif layout in ('code', 'code-block', 'snippet', 'codeblock'):
            # Bloc de code : encadré sombre avec liséré accent à gauche.
            # Champs tolérés : 'code' (le source), 'language'/'lang'
            # (étiquette), 'description'/'intro' (phrase au-dessus).
            code = str(contenu.get('code') or contenu.get('source') or '')
            language = str(contenu.get('language') or contenu.get('lang') or '')
            description = str(contenu.get('description') or contenu.get('intro')
                              or contenu.get('explication') or '')
            if code.strip() and _make_code_slide(prs, titre, code, language,
                                                  description=description) is None:
                # Fallback : on rend le code en bloc texte mono via content_slide
                _make_content_slide(prs, titre, code.split('\n')[:12],
                                    is_bullets=False)

        elif layout in ('case', 'case-study', 'case_study', 'cas-pratique',
                        'cas_pratique', 'cas-pratique', 'etude-cas', 'etude_de_cas'):
            # Cas pratique : 4 cartes Contexte / Problème / Solution / Résultat.
            # Champs tolérés (souples) — chaque case accepte plusieurs alias.
            contexte = str(contenu.get('contexte') or contenu.get('context')
                           or contenu.get('situation') or contenu.get('cadre') or '')
            probleme = str(contenu.get('probleme') or contenu.get('problème')
                           or contenu.get('problem') or contenu.get('enjeu')
                           or contenu.get('defi') or contenu.get('défi') or '')
            solution = str(contenu.get('solution') or contenu.get('approche')
                           or contenu.get('action') or contenu.get('strategie')
                           or contenu.get('stratégie') or '')
            resultat = str(contenu.get('resultat') or contenu.get('résultat')
                           or contenu.get('outcome') or contenu.get('impact')
                           or contenu.get('bilan') or '')
            if _make_case_study_slide(prs, titre, contexte, probleme,
                                      solution, resultat) is None:
                # Fallback : bullets simples avec labels en gras
                items = [f'**{lbl}** : {body}' for lbl, body in (
                    ('Contexte', contexte),
                    ('Problème', probleme),
                    ('Solution', solution),
                    ('Résultat', resultat),
                ) if body and body.strip()]
                if items:
                    _make_content_slide(prs, titre or 'Cas pratique',
                                        items, is_bullets=True)

        elif layout in ('timeline', 'chronologie', 'frise', 'history', 'historique'):
            # Frise chronologique : 2 à 7 événements datés.
            # Champs tolérés pour la liste : 'events', 'evenements', 'dates',
            # 'jalons', 'items'. Chaque entrée peut être :
            #   {'date': '1995', 'titre': '...', 'description': '...'}
            #   ou un simple string (sera affiché comme titre seul).
            raw = (contenu.get('events') or contenu.get('evenements')
                   or contenu.get('événements') or contenu.get('dates')
                   or contenu.get('jalons') or contenu.get('items') or [])
            if raw and _make_timeline_slide(prs, titre, raw) is None:
                # Fallback : bullets simples si la frise n'a pas pu être rendue
                texts = []
                for e in raw:
                    if isinstance(e, dict):
                        d = str(e.get('date') or '')
                        t = str(e.get('titre') or e.get('title') or '')
                        texts.append(f'{d} — {t}' if d and t else (t or d))
                    else:
                        texts.append(str(e))
                texts = [t for t in texts if t]
                if texts:
                    _make_content_slide(prs, titre, texts, is_bullets=True)

        elif layout in ('versus', 'comparison', 'compare', 'comparaison', 'vs'):
            # Versus / comparaison : deux colonnes avec en-têtes colorés
            # (vert/orange si sémantique pro/con, accent violet sinon).
            # Champs tolérés (souples) :
            #   {'left_label', 'left_items', 'right_label', 'right_items'}
            #   ou {'avant': {...}, 'apres': {...}} (avec label+items)
            ll = (contenu.get('left_label') or contenu.get('gauche')
                  or (contenu.get('avant') or {}).get('label') or 'Avant')
            li = (contenu.get('left_items') or contenu.get('items_gauche')
                  or (contenu.get('avant') or {}).get('items') or [])
            rl = (contenu.get('right_label') or contenu.get('droite')
                  or (contenu.get('apres') or {}).get('label') or 'Après')
            ri = (contenu.get('right_items') or contenu.get('items_droite')
                  or (contenu.get('apres') or {}).get('items') or [])
            li = [str(x) for x in li if str(x).strip()]
            ri = [str(x) for x in ri if str(x).strip()]
            if li and ri:
                if _make_versus_slide(prs, titre, str(ll), li, str(rl), ri) is None:
                    _make_two_column_slide(prs, titre, li + ri)

        elif layout in ('objectives', 'objectifs', 'goals', 'learning-objectives'):
            # Objectifs pédagogiques : liste de compétences visées (3 à 6).
            # Champs tolérés : 'objectifs', 'items', 'goals', 'competences'.
            raw = (contenu.get('objectifs') or contenu.get('items')
                   or contenu.get('goals') or contenu.get('competences') or [])
            objs = [str(x) for x in raw if str(x).strip()]
            if objs and _make_objectives_slide(prs, objs) is None:
                _make_content_slide(prs, titre or 'Objectifs pédagogiques',
                                    objs, is_bullets=True)

        elif layout in ('synthese', 'synthèse', 'recap', 'summary', 'takeaways'):
            # Synthèse / À retenir : 3 à 5 points clés récapitulatifs.
            # Champs tolérés : 'points', 'items', 'takeaways', 'recap'.
            raw = (contenu.get('points') or contenu.get('items')
                   or contenu.get('takeaways') or contenu.get('recap') or [])
            pts = [str(x) for x in raw if str(x).strip()]
            if pts and _make_synthese_slide(prs, pts) is None:
                _make_content_slide(prs, titre or 'Synthèse', pts,
                                    is_bullets=True)

        elif layout in ('stepper', 'steps', 'process', 'processus', 'etapes', 'étapes'):
            # Stepper vertical : processus séquentiel numéroté.
            # Champs tolérés pour la liste : 'etapes', 'steps', 'phases', 'items'.
            # Chaque étape peut être un dict {titre, description} ou un simple string.
            raw_steps = (contenu.get('etapes') or contenu.get('étapes')
                         or contenu.get('steps') or contenu.get('phases')
                         or contenu.get('items') or [])
            if raw_steps and _make_stepper_slide(prs, titre, raw_steps) is None:
                # Fallback si moins de 2 étapes : bullets simple
                texts = []
                for s in raw_steps:
                    if isinstance(s, dict):
                        t = s.get('titre') or s.get('title') or ''
                        d = s.get('description') or s.get('desc') or ''
                        texts.append(f'{t} : {d}' if t and d else (t or d))
                    else:
                        texts.append(str(s))
                texts = [t for t in texts if t]
                if texts:
                    _make_content_slide(prs, titre, texts, is_bullets=True)

        else:
            # Fallback : bullets avec les valeurs du contenu
            fallback = [str(v) for v in contenu.values() if v][:8]
            if fallback:
                _make_content_slide(prs, titre, fallback, is_bullets=True)

    # Remplit ou supprime la slide sommaire selon le nombre de sections captées.
    # Si on supprime la TOC, les indices section décalent vers le haut de 1.
    toc_removed = False
    if toc_slide is not None:
        if len(toc_entries) >= 2:
            _fill_toc_slide(toc_slide, toc_entries)
        else:
            _remove_slide(prs, toc_slide)
            toc_removed = True

    if toc_removed:
        # La TOC était à l'index 1 ; chaque slide après recule de 1
        section_slide_indices = {i - 1 if i > 1 else i for i in section_slide_indices}

    _apply_footers(prs, specialite, module, chapitre,
                   skip_indices=section_slide_indices)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ═══════════════════════════════════════════════════════
#  POINT D'ENTRÉE
# ═══════════════════════════════════════════════════════

def markdown_to_pptx(contenu: str, specialite: str, module: str,
                     chapitre: str, niveau: str = '',
                     title_image: bytes = None, photographer: str = '') -> bytes:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _make_title_slide(prs, specialite, module, chapitre, niveau,
                      title_image=title_image, photographer=photographer)

    # Sommaire cliquable : créé comme placeholder juste après la titre, rempli
    # à la fin une fois que toutes les slides de section ont été créées.
    toc_slide = _make_toc_placeholder(prs)
    toc_entries = []  # [(titre_section, première_slide_de_la_section), ...]

    SKIP      = {'tableau comparatif', 'synthèse visuelle', 'pour aller plus loin'}
    SEC_DEF   = {'définitions des concepts clés', 'définitions'}
    SEC_PTS   = {'points importants à retenir', 'points clés'}
    # Mots-clés déclenchant une slide "callout" (citation / règle clé / maxime)
    SEC_CALLOUT = {'citation', 'règle clé', 'principe clé', 'principe fondamental',
                   'maxime', 'à méditer'}
    # Mots-clés déclenchant une slide "stepper" (processus séquentiel)
    SEC_STEPPER = {'étapes', 'etapes', 'procédure', 'procedure',
                   'démarche', 'demarche', 'processus', 'méthodologie',
                   'methodologie', 'phases', 'étapes clés', 'cycle de vie'}
    # Mots-clés déclenchant une slide "KPI / chiffres clés"
    SEC_KPI = {'chiffres clés', 'chiffres cles', 'kpi', 'indicateurs clés',
               'indicateurs cles', 'statistiques clés', 'statistiques cles',
               'quelques chiffres', 'en chiffres', 'données chiffrées',
               'donnees chiffrees', 'le marché en chiffres',
               'le marche en chiffres'}
    # Mots-clés déclenchant une slide "code source" (mono + liséré accent).
    # Couvre aussi l'auto-détection : si le corps de la section contient un
    # bloc ``` et que ce bloc fait >40% du texte, on déclenche aussi.
    SEC_CODE = {'code', 'extrait de code', 'extrait code', 'snippet',
                'exemple de code', 'exemple code', 'implémentation',
                'implementation', 'pseudocode', 'pseudo-code',
                'pseudo code', 'syntaxe', 'illustration code',
                'requête sql', 'requete sql', 'script', 'commande'}
    # Mots-clés déclenchant une slide "cas pratique" (4 cartes)
    SEC_CASE = {'cas pratique', 'cas pratiques', 'étude de cas', 'etude de cas',
                'études de cas', 'etudes de cas', 'business case',
                'cas concret', 'cas d\'école', 'cas decole',
                'cas d\'application', 'mise en situation',
                'illustration concrète', 'illustration concrete',
                'case study', 'exemple détaillé', 'exemple detaille',
                'situation réelle', 'situation reelle'}
    # Mots-clés déclenchant une slide "frise chronologique"
    SEC_TIMELINE = {'chronologie', 'timeline', 'frise', 'frise chronologique',
                    'historique', 'évolution historique', 'evolution historique',
                    'dates clés', 'dates cles', 'jalons', 'jalons historiques',
                    'étapes historiques', 'etapes historiques',
                    'histoire de', 'genèse', 'genese'}
    # Mots-clés déclenchant une slide "versus / comparaison"
    SEC_VERSUS = {'avant / après', 'avant/après', 'avant apres', 'avant-après',
                  'avant-apres', ' versus ', 'versus', ' vs ',
                  'pour et contre', 'avantages et inconvénients',
                  'avantages et inconvenients', 'forces et faiblesses',
                  'traditionnel vs moderne', 'traditionnel / moderne',
                  'ancien vs nouveau', 'théorie vs pratique',
                  'theorie vs pratique'}
    # Mots-clés déclenchant une slide "Objectifs pédagogiques" (tête de chapitre).
    # Volontairement spécifiques : 'objectifs' bare matcherait « Objectifs
    # commerciaux / marketing » dans le corps d'un cours — à proscrire.
    SEC_GOAL = {'objectifs pédagogiques', 'objectifs pedagogiques',
                'objectifs du cours', 'objectifs du chapitre',
                'objectifs de la séquence', 'objectifs de la sequence',
                'objectifs d\'apprentissage', 'objectifs dapprentissage',
                'compétences visées', 'competences visees',
                'à l’issue', 'à l\'issue', 'a l\'issue'}
    # Mots-clés déclenchant une slide "Synthèse / À retenir" (fin de chapitre)
    SEC_RECAP = {'synthèse', 'synthese', 'à retenir', 'a retenir',
                 'récapitulatif', 'recapitulatif', 'take-aways', 'take aways',
                 'en résumé', 'en resume', 'conclusion clé', 'l\'essentiel',
                 'l’essentiel'}

    section_counter = 0
    section_slide_indices = set()  # indices à exclure du footer (slides full-bleed)

    for section in re.split(r'\n(?=## )', contenu):
        section = section.strip()
        if not section.startswith('##'):
            continue

        lines        = section.split('\n')
        h2_title     = _clean(lines[0].lstrip('#').strip())
        section_body = '\n'.join(lines[1:])
        h2_lower     = h2_title.lower()

        if any(kw in h2_lower for kw in SKIP):
            continue

        # Snapshot avant rendu : permet ensuite de capter la 1re slide générée
        # pour cette section et de l'ajouter au sommaire cliquable.
        _slides_before = len(prs.slides)

        def _record_toc():
            """Ajoute l'entrée au sommaire si au moins une slide a été créée."""
            if len(prs.slides) > _slides_before:
                toc_entries.append((h2_title, prs.slides[_slides_before]))

        # ── Objectifs pédagogiques (tête de chapitre/section) ──
        if any(kw in h2_lower for kw in SEC_GOAL):
            objs = _extract_bullets(section_body, max_b=6)
            if not objs:
                # Fallback : lignes non vides du 1er paragraphe
                para = _first_paragraph(section_body)
                if para:
                    objs = [l.strip() for l in para.split('.') if l.strip()][:6]
            if objs and _make_objectives_slide(prs, objs) is not None:
                _record_toc()
                continue
            # Sinon : fallback sur dispatch classique ci-dessous

        # ── Synthèse / À retenir (fin de chapitre/section) ──
        if any(kw in h2_lower for kw in SEC_RECAP):
            pts = _extract_bullets(section_body, max_b=5)
            if not pts:
                para = _first_paragraph(section_body)
                if para:
                    pts = [l.strip() for l in para.split('.') if l.strip()][:5]
            if pts and _make_synthese_slide(prs, pts) is not None:
                _record_toc()
                continue

        # ── Frise chronologique / timeline ─────────────────────────────
        # Déclenchement : mot-clé titre OU >=3 bullets reconnus comme dates.
        # Placé avant KPI pour que des bullets « 1995 — ... » ne soient pas
        # mépris pour des chiffres clés (la date est numérique).
        if any(kw in h2_lower for kw in SEC_TIMELINE) or _looks_like_timeline(section_body):
            events = _parse_timeline(section_body, max_n=7)
            if len(events) >= 2:
                if _make_timeline_slide(prs, h2_title, events) is not None:
                    _record_toc()
                    continue

        # ── KPI / chiffres clés ───────────────────────────────────────
        # Déclenchement : mot-clé titre OU auto-détection (>=50% des bullets
        # sont des pourcentages/montants). Cascade vers progress-bars si tous
        # en %, bar-chart si au moins 2 numériques, sinon cartes texte.
        if any(kw in h2_lower for kw in SEC_KPI) or _looks_like_kpi(section_body):
            kpis = _parse_kpi_bullets(section_body, max_n=4)
            if len(kpis) >= 2:
                # Convertit au format attendu par _make_stat_slide et cie
                stats = [{'valeur': k['valeur'], 'label': k['label']} for k in kpis]
                if _make_progress_slide(prs, h2_title, stats) is None:
                    if _make_stat_chart_slide(prs, h2_title, stats) is None:
                        _make_stat_slide(prs, h2_title, stats)
                _record_toc()
                continue
            # Sinon on retombe sur le dispatch classique

        if any(kw in h2_lower for kw in SEC_DEF):
            items = _parse_definitions(section_body)
            if items:
                _make_definitions_slide(prs, items)
            _record_toc()
            continue

        if any(kw in h2_lower for kw in SEC_PTS):
            points = _extract_bullets(section_body, max_b=12)
            if not points:
                points = _wrap_text(_first_paragraph(section_body), 120)
            if points:
                _make_key_points_slide(prs, points)
            _record_toc()
            continue

        if any(kw in h2_lower for kw in SEC_CALLOUT):
            # D'abord un blockquote, sinon on prend le premier paragraphe.
            quote, attribution = _parse_blockquote(section_body)
            if not quote:
                quote = _first_paragraph(section_body)
            if quote:
                _make_callout_slide(prs, quote, attribution)
            _record_toc()
            continue

        # Auto-détection : section dont le corps est exclusivement un blockquote
        # (titre non matché par SEC_CALLOUT mais contenu sans ambiguïté)
        if _is_pure_blockquote(section_body):
            quote, attribution = _parse_blockquote(section_body)
            if quote:
                _make_callout_slide(prs, quote, attribution, section_label=h2_title)
                _record_toc()
                continue

        # ── Bloc de code source ───────────────────────────────────────
        # Déclenchement : mot-clé titre OU corps contenant un bloc ``` qui
        # represente >40% des caractères du body (auto-détection : la
        # section est centrée sur le code, pas sur du texte explicatif).
        cb = _parse_code_blocks(section_body)
        cb_dominant = (cb is not None
                       and len(cb['code']) > 0.4 * max(1, len(section_body)))
        if any(kw in h2_lower for kw in SEC_CODE) or cb_dominant:
            if cb and cb['code'].strip():
                if _make_code_slide(
                    prs, h2_title, cb['code'], cb['language'],
                    description=cb['before'],
                ) is not None:
                    _record_toc()
                    continue
            # Sinon : pas de bloc de code malgré le mot-clé → fallback

        # ── Cas pratique / étude de cas (4 cartes Z-flow) ─────────────
        # Déclenchement : mot-clé titre. Structure du corps détectée par
        # `_parse_case_study` (H3 / labels gras / bullets gras).
        if any(kw in h2_lower for kw in SEC_CASE):
            case = _parse_case_study(section_body)
            if case and _make_case_study_slide(
                prs, h2_title,
                case['contexte'], case['probleme'],
                case['solution'], case['resultat'],
            ) is not None:
                _record_toc()
                continue
            # Sinon : fallback sur le dispatch classique

        # ── Versus / comparaison (2 colonnes sémantiques) ─────────────
        # Déclenchement : mot-clé titre. Structure du corps détectée par
        # `_parse_versus` (2 H3 / 2-col table / pattern inline).
        if any(kw in h2_lower for kw in SEC_VERSUS):
            parsed = _parse_versus(section_body)
            if parsed and _make_versus_slide(
                prs, h2_title,
                parsed['left_label'], parsed['left_items'],
                parsed['right_label'], parsed['right_items'],
            ) is not None:
                _record_toc()
                continue
            # Sinon : fallback sur le dispatch classique

        # Stepper : titre déclencheur OU contenu principalement en liste numérotée
        if any(kw in h2_lower for kw in SEC_STEPPER) or _looks_like_steps(section_body):
            steps = _parse_steps(section_body)
            if len(steps) >= 2:
                if _make_stepper_slide(prs, h2_title, steps) is not None:
                    _record_toc()
                    continue
            # Sinon, on retombe sur le dispatch classique ci-dessous

        section_counter += 1
        _make_section_slide(prs, h2_title, numero=section_counter)
        section_slide_indices.add(len(prs.slides) - 1)

        subsections    = [s.strip() for s in re.split(r'\n(?=### )', section_body) if s.strip()]
        has_subsections = any(s.startswith('###') for s in subsections)

        if has_subsections:
            intro = subsections[0] if not subsections[0].startswith('###') else ''
            if intro:
                headers, rows = _parse_md_table(intro)
                if headers and rows:
                    _make_table_slide(prs, h2_title, headers, rows, section_label=h2_title)
                else:
                    para = _first_paragraph(intro)
                    if para:
                        _make_content_slide(prs, h2_title, _wrap_text(para, 95),
                                            is_bullets=False, section_label=h2_title)

            for sub in subsections:
                if not sub.startswith('###'):
                    continue
                sub_lines = sub.split('\n')
                h3_title  = _clean(sub_lines[0].lstrip('#').strip())
                sub_body  = '\n'.join(sub_lines[1:])

                headers, rows = _parse_md_table(sub_body)
                if headers and rows:
                    _make_table_slide(prs, h3_title, headers, rows, section_label=h2_title)
                    continue

                bullets = _extract_bullets(sub_body)
                if bullets:
                    if len(bullets) > MAX_BULLETS_SINGLE:
                        _make_two_column_slide(prs, h3_title, bullets[:MAX_BULLETS_COL * 2],
                                               section_label=h2_title)
                    else:
                        _make_content_slide(prs, h3_title, bullets,
                                            is_bullets=True, section_label=h2_title)
                else:
                    para = _first_paragraph(sub_body)
                    if para:
                        _make_content_slide(prs, h3_title, _wrap_text(para, 95),
                                            is_bullets=False, section_label=h2_title)
        else:
            headers, rows = _parse_md_table(section_body)
            if headers and rows:
                _make_table_slide(prs, h2_title, headers, rows)
            else:
                bullets = _extract_bullets(section_body)
                if bullets:
                    if len(bullets) > MAX_BULLETS_SINGLE:
                        _make_two_column_slide(prs, h2_title, bullets[:MAX_BULLETS_COL * 2])
                    else:
                        _make_content_slide(prs, h2_title, bullets, is_bullets=True)
                else:
                    para = _first_paragraph(section_body)
                    if para:
                        _make_content_slide(prs, h2_title, _wrap_text(para, 95), is_bullets=False)

        # Fin du bloc par fall-through (section slide + subsections / bullets)
        _record_toc()

    # Remplit ou supprime la slide sommaire selon les entrées collectées.
    # Si on supprime la TOC (idx 1), tous les indices > 1 reculent de 1.
    toc_removed = False
    if len(toc_entries) >= 2:
        _fill_toc_slide(toc_slide, toc_entries)
    else:
        _remove_slide(prs, toc_slide)
        toc_removed = True

    if toc_removed:
        section_slide_indices = {i - 1 if i > 1 else i for i in section_slide_indices}

    _apply_footers(prs, specialite, module, chapitre,
                   skip_indices=section_slide_indices)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
